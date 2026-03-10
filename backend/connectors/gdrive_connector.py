"""
Google Drive Connector for 2nd Brain
Syncs documents and files from Google Drive
"""

import os
import io
from datetime import datetime, timezone
from typing import List, Optional, Dict, Any
from urllib.parse import urlencode
from concurrent.futures import ThreadPoolExecutor, as_completed
import requests

from connectors.base_connector import (
    BaseConnector,
    ConnectorConfig,
    ConnectorStatus,
    Document
)

try:
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    GDRIVE_AVAILABLE = True
except ImportError:
    GDRIVE_AVAILABLE = False
    print("[GDrive] Google API SDK not installed. Run: pip install google-api-python-client google-auth")


# Google Drive scopes
GDRIVE_SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
]

# Google Workspace MIME types → export format for content extraction
GOOGLE_WORKSPACE_TYPES = {
    'application/vnd.google-apps.document': 'text/plain',
    'application/vnd.google-apps.spreadsheet': 'text/csv',
    'application/vnd.google-apps.presentation': 'text/plain',
    'application/vnd.google-apps.drawing': 'image/png',
    'application/vnd.google-apps.form': 'text/plain',
}

# Non-Workspace file types we can extract content from
EXTRACTABLE_TYPES = {
    'application/pdf': None,
    'text/plain': None,
    'text/markdown': None,
    'text/html': None,
    'text/csv': None,
    'text/xml': None,
    'application/json': None,
    'application/xml': None,
    'application/rtf': None,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': None,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': None,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': None,
    'application/msword': None,
    'application/vnd.ms-excel': None,
    'application/vnd.ms-powerpoint': None,
    'image/png': None,
    'image/jpeg': None,
    'image/gif': None,
    'image/webp': None,
    'image/tiff': None,
    'video/mp4': None,
    'video/quicktime': None,
    'audio/mpeg': None,
    'audio/wav': None,
}

# MIME types to EXCLUDE from file listing (non-file objects)
EXCLUDED_TYPES = {
    'application/vnd.google-apps.folder',
    'application/vnd.google-apps.shortcut',
    'application/vnd.google-apps.fusiontable',
    'application/vnd.google-apps.map',
    'application/vnd.google-apps.site',
}


class GDriveConnector(BaseConnector):
    """Connector for Google Drive - synchronous implementation"""

    CONNECTOR_TYPE = "gdrive"
    REQUIRED_CREDENTIALS = ["access_token"]
    MAX_CONCURRENT_FILES = 20
    OPTIONAL_SETTINGS = {
        "folder_ids": [],              # Specific folders to sync (empty = all)
        "include_shared": True,        # Include shared files
        "include_trashed": False,      # Include trashed files
        "max_file_size_mb": 25,        # Skip files larger than this
        "max_files": 500               # Maximum files to sync
    }

    def __init__(self, config: ConnectorConfig):
        super().__init__(config)
        self.service = None
        self.credentials: Optional[Credentials] = None
        self.user_email: str = ""

    def connect(self) -> bool:
        """Establish connection to Google Drive API"""
        if not GDRIVE_AVAILABLE:
            self._set_error("Google API SDK not installed")
            return False

        try:
            self.status = ConnectorStatus.CONNECTING

            access_token = self.config.credentials.get("access_token")
            refresh_token = self.config.credentials.get("refresh_token")

            if not access_token:
                self._set_error("Missing access_token")
                return False

            self.credentials = Credentials(
                token=access_token,
                refresh_token=refresh_token,
                token_uri="https://oauth2.googleapis.com/token",
                client_id=os.getenv("GOOGLE_CLIENT_ID"),
                client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
                scopes=GDRIVE_SCOPES
            )

            # Refresh if expired
            if self.credentials.expired and self.credentials.refresh_token:
                try:
                    self.credentials.refresh(Request())
                    self.config.credentials["access_token"] = self.credentials.token
                except Exception as e:
                    print(f"[GDrive] Token refresh failed: {e}")

            self.service = build('drive', 'v3', credentials=self.credentials)

            # Test connection
            about = self.service.about().get(fields="user").execute()
            self.user_email = about.get("user", {}).get("emailAddress", "")

            self.status = ConnectorStatus.CONNECTED
            self._clear_error()
            print(f"[GDrive] Connected as: {self.user_email}")
            return True

        except Exception as e:
            self._set_error(f"Connection failed: {str(e)}")
            return False

    def disconnect(self) -> bool:
        """Disconnect from Google Drive"""
        self.service = None
        self.credentials = None
        self.status = ConnectorStatus.DISCONNECTED
        return True

    def test_connection(self) -> bool:
        """Verify connection is still valid"""
        if not self.service:
            return False
        try:
            self.service.about().get(fields="user").execute()
            return True
        except Exception:
            return False

    def sync(self, since: Optional[datetime] = None) -> List[Document]:
        """Sync files from Google Drive"""
        if not self.service:
            if not self.connect():
                return []

        self.status = ConnectorStatus.SYNCING
        documents = []

        try:
            max_files = self.config.settings.get("max_files", 500)
            include_shared = self.config.settings.get("include_shared", True)
            include_trashed = self.config.settings.get("include_trashed", False)
            max_size = self.config.settings.get("max_file_size_mb", 25) * 1024 * 1024
            folder_ids = self.config.settings.get("folder_ids", [])

            print(f"[GDrive] Starting sync (max_files={max_files})")

            # Build query
            query_parts = []

            # Exclude folders, shortcuts, and other non-file types
            for excluded_type in EXCLUDED_TYPES:
                query_parts.append(f"mimeType!='{excluded_type}'")

            if not include_trashed:
                query_parts.append("trashed=false")

            if since:
                since_str = since.strftime("%Y-%m-%dT%H:%M:%S")
                query_parts.append(f"modifiedTime > '{since_str}'")

            if folder_ids:
                folder_queries = [f"'{fid}' in parents" for fid in folder_ids]
                query_parts.append(f"({' or '.join(folder_queries)})")

            query = " and ".join(query_parts)

            # STEP 1: Collect all eligible file metadata
            file_items = []
            page_token = None

            while len(file_items) < max_files:
                response = self.service.files().list(
                    q=query,
                    spaces='drive',
                    fields='nextPageToken, files(id, name, mimeType, size, modifiedTime, createdTime, owners, webViewLink, parents)',
                    pageToken=page_token,
                    pageSize=min(100, max_files - len(file_items)),
                    includeItemsFromAllDrives=include_shared,
                    supportsAllDrives=True
                ).execute()

                files = response.get('files', [])
                print(f"[GDrive] Found {len(files)} files in batch")

                for file in files:
                    size = int(file.get('size', 0) or 0)
                    if size > max_size:
                        continue
                    file_items.append(file)
                    if len(file_items) >= max_files:
                        break

                page_token = response.get('nextPageToken')
                if not page_token:
                    break

            print(f"[GDrive] Collected {len(file_items)} files, processing in parallel (max {self.MAX_CONCURRENT_FILES} concurrent)...")

            # STEP 2: Process all files in PARALLEL
            file_count = 0
            with ThreadPoolExecutor(max_workers=self.MAX_CONCURRENT_FILES) as executor:
                futures = {executor.submit(self._file_to_document, f): f for f in file_items}
                for future in as_completed(futures):
                    try:
                        doc = future.result()
                        if doc:
                            documents.append(doc)
                            file_count += 1
                            if self.on_document_ready:
                                try:
                                    self.on_document_ready(doc)
                                except Exception as cb_err:
                                    print(f"[GDrive] on_document_ready error: {cb_err}")
                            print(f"[GDrive] Processed file {file_count}/{len(file_items)}: {doc.title[:50] if doc.title else 'Untitled'}")
                    except Exception as e:
                        print(f"[GDrive] Error processing file: {e}")

            self.config.last_sync = datetime.now(timezone.utc)
            self.status = ConnectorStatus.CONNECTED
            print(f"[GDrive] Sync complete: {len(documents)} files")

        except Exception as e:
            self._set_error(f"Sync failed: {str(e)}")
            import traceback
            traceback.print_exc()

        return documents

    def get_document(self, doc_id: str) -> Optional[Document]:
        """Retrieve a specific file by ID"""
        if not self.service:
            self.connect()

        try:
            file_id = doc_id.replace("gdrive_", "")
            file = self.service.files().get(
                fileId=file_id,
                fields='id, name, mimeType, size, modifiedTime, createdTime, owners, webViewLink, parents'
            ).execute()
            return self._file_to_document(file)
        except Exception as e:
            self._set_error(f"Failed to get document: {str(e)}")
            return None

    def _file_to_document(self, file: Dict[str, Any]) -> Optional[Document]:
        """Convert Google Drive file to Document object"""
        try:
            file_id = file['id']
            name = file['name']
            mime_type = file['mimeType']

            content = self._extract_content(file_id, mime_type, name)

            if not content:
                # Still include the file with a placeholder so it appears in the selection modal
                content = f"[File: {name}] (Content extraction not available for type: {mime_type})"
                print(f"[GDrive] No content extracted for {name} ({mime_type}), using placeholder", flush=True)

            owners = file.get('owners', [])
            owner_names = [o.get('displayName', o.get('emailAddress', '')) for o in owners]

            metadata = {
                "file_id": file_id,
                "mime_type": mime_type,
                "size": file.get('size'),
                "owners": owner_names,
                "user": self.user_email
            }

            modified = file.get('modifiedTime', '')
            timestamp = None
            if modified:
                try:
                    timestamp = datetime.fromisoformat(modified.replace('Z', '+00:00'))
                except:
                    pass

            return Document(
                doc_id=f"gdrive_{file_id}",
                source="gdrive",
                content=content,
                title=name,
                metadata=metadata,
                timestamp=timestamp,
                author=owner_names[0] if owner_names else None,
                url=file.get('webViewLink'),
                doc_type="file"
            )

        except Exception as e:
            print(f"[GDrive] Error converting file: {e}")
            return None

    def _get_file_type(self, mime_type: str) -> str:
        """Get human-readable file type"""
        if 'document' in mime_type or 'wordprocessing' in mime_type:
            return "Document"
        elif 'spreadsheet' in mime_type:
            return "Spreadsheet"
        elif 'presentation' in mime_type:
            return "Presentation"
        elif 'pdf' in mime_type:
            return "PDF"
        else:
            return "File"

    def _extract_content(self, file_id: str, mime_type: str, filename: str) -> Optional[str]:
        """Extract text content from file"""
        try:
            # Google Workspace files - export as text
            if mime_type.startswith('application/vnd.google-apps.'):
                export_type = GOOGLE_WORKSPACE_TYPES.get(mime_type)
                if export_type:
                    response = self.service.files().export(
                        fileId=file_id,
                        mimeType=export_type
                    ).execute()

                    if isinstance(response, bytes):
                        return response.decode('utf-8', errors='ignore')
                    return str(response)
                return None

            # Regular files - download and extract
            request = self.service.files().get_media(fileId=file_id)
            file_buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(file_buffer, request)

            done = False
            while not done:
                _, done = downloader.next_chunk()

            file_content = file_buffer.getvalue()

            # Handle based on type
            # Plain text types (txt, md, json, csv, xml, rtf)
            if mime_type in ['text/plain', 'text/markdown', 'application/json', 'text/csv',
                             'text/xml', 'application/xml', 'application/rtf']:
                return file_content.decode('utf-8', errors='ignore')

            elif mime_type == 'text/html':
                try:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(file_content, 'html.parser')
                    return soup.get_text(separator='\n')
                except ImportError:
                    import re
                    text = file_content.decode('utf-8', errors='ignore')
                    return re.sub('<[^<]+?>', '', text)

            elif (mime_type in ('application/pdf',) or
                  'wordprocessingml' in mime_type or 'spreadsheetml' in mime_type or
                  'presentationml' in mime_type or
                  mime_type in ('application/msword', 'application/vnd.ms-excel', 'application/vnd.ms-powerpoint')):
                # Use Mistral Document AI (with local fallback) for PDF, DOCX, XLSX, PPTX, DOC, XLS, PPT
                try:
                    from parsers.document_parser import DocumentParser
                    parser = DocumentParser()
                    parsed = parser.parse_file_bytes(file_content, filename)
                    if parsed:
                        return parsed
                except Exception as e:
                    print(f"[GDrive] Mistral parse failed for {filename}: {e}")
                # Fallback for DOCX
                if 'wordprocessingml' in mime_type or mime_type == 'application/msword':
                    try:
                        import docx
                        doc = docx.Document(io.BytesIO(file_content))
                        return '\n'.join([p.text for p in doc.paragraphs])
                    except Exception:
                        pass
                # Fallback for XLSX
                if 'spreadsheetml' in mime_type or mime_type == 'application/vnd.ms-excel':
                    try:
                        import pandas as pd
                        df = pd.read_excel(io.BytesIO(file_content))
                        return df.to_string()
                    except Exception:
                        pass
                # Fallback for PPTX
                if 'presentationml' in mime_type:
                    try:
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(file_content))
                        text_parts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_parts.append(shape.text)
                        return '\n'.join(text_parts)
                    except Exception:
                        pass
                return None

            elif mime_type and mime_type.startswith('image/'):
                # OCR via Azure Document Intelligence
                try:
                    from parsers.document_parser import DocumentParser
                    parser = DocumentParser()
                    parsed = parser.parse_file_bytes(file_content, filename)
                    if parsed:
                        print(f"[GDrive] Image OCR extracted {len(parsed)} chars from {filename}")
                        return parsed
                except Exception as e:
                    print(f"[GDrive] Image OCR failed for {filename}: {e}")
                return None

            elif mime_type and (mime_type.startswith('video/') or mime_type.startswith('audio/')):
                # Whisper transcription for audio and video
                try:
                    from services.knowledge_service import KnowledgeService
                    from database.models import SessionLocal
                    db = SessionLocal()
                    try:
                        ks = KnowledgeService(db)
                        result = ks.transcribe_audio(file_content, filename)
                        if result and result.text:
                            print(f"[GDrive] Transcribed {len(result.text)} chars from {filename}")
                            return result.text
                    finally:
                        db.close()
                except Exception as e:
                    print(f"[GDrive] Transcription failed for {filename}: {e}")
                return None

            return None

        except Exception as e:
            print(f"[GDrive] Content extraction error for {filename}: {e}")
            return None

    @classmethod
    def get_auth_url(cls, redirect_uri: str, state: str) -> str:
        """Generate Google OAuth authorization URL"""
        client_id = os.getenv("GOOGLE_CLIENT_ID", "")

        params = {
            'client_id': client_id,
            'redirect_uri': redirect_uri,
            'scope': ' '.join(GDRIVE_SCOPES),
            'response_type': 'code',
            'access_type': 'offline',
            'prompt': 'consent',
            'state': state
        }

        return f"https://accounts.google.com/o/oauth2/v2/auth?{urlencode(params)}"

    @classmethod
    def exchange_code(cls, code: str, redirect_uri: str) -> Dict[str, Any]:
        """Exchange authorization code for tokens"""
        response = requests.post(
            "https://oauth2.googleapis.com/token",
            data={
                'client_id': os.getenv("GOOGLE_CLIENT_ID"),
                'client_secret': os.getenv("GOOGLE_CLIENT_SECRET"),
                'code': code,
                'grant_type': 'authorization_code',
                'redirect_uri': redirect_uri
            }
        )

        data = response.json()

        if 'error' in data:
            raise Exception(f"OAuth failed: {data.get('error_description', data['error'])}")

        return {
            'access_token': data['access_token'],
            'refresh_token': data.get('refresh_token'),
            'expires_in': data.get('expires_in'),
            'token_type': data.get('token_type')
        }
