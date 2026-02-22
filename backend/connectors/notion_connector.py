"""
Notion Connector for 2nd Brain
Syncs pages and databases from Notion workspaces
"""

import os
import io
import base64
import tempfile
from datetime import datetime, timezone
from typing import List, Optional, Dict, Any
import requests

from connectors.base_connector import (
    BaseConnector,
    ConnectorConfig,
    ConnectorStatus,
    Document
)

try:
    from notion_client import Client
    NOTION_AVAILABLE = True
except ImportError:
    NOTION_AVAILABLE = False
    print("[Notion] notion-client not installed. Run: pip install notion-client")


class NotionConnector(BaseConnector):
    """Connector for Notion workspaces - synchronous implementation"""

    CONNECTOR_TYPE = "notion"
    REQUIRED_CREDENTIALS = ["access_token"]
    OPTIONAL_SETTINGS = {
        "database_ids": [],           # Specific databases to sync (empty = all)
        "include_archived": False,    # Include archived pages
        "max_blocks_per_page": 1000,  # Limit blocks fetched per page
        "max_pages": 500,             # Maximum pages to sync
        "max_depth": 5                # Max nesting depth for child blocks
    }

    def __init__(self, config: ConnectorConfig):
        super().__init__(config)
        self.client: Optional[Client] = None
        self.workspace_name: str = ""
        self.workspace_id: str = ""
        # Accumulates separate Document objects for files embedded in pages
        self._child_documents: List[Document] = []
        self._current_page_id: str = ""
        self._current_page_url: str = ""
        self._current_page_timestamp: Optional[datetime] = None

    def connect(self) -> bool:
        """Establish connection to Notion API"""
        if not NOTION_AVAILABLE:
            self._set_error("Notion SDK not installed")
            return False

        try:
            self.status = ConnectorStatus.CONNECTING
            access_token = self.config.credentials.get("access_token")

            if not access_token:
                self._set_error("Missing access_token")
                return False

            self.client = Client(auth=access_token)

            # Test connection and get workspace info
            me = self.client.users.me()
            self.workspace_name = me.get("name", "Notion Workspace")

            bot_info = me.get("bot", {})
            owner = bot_info.get("owner", {})
            if owner.get("type") == "workspace":
                self.workspace_id = owner.get("workspace", {}).get("id", "")

            self.status = ConnectorStatus.CONNECTED
            self._clear_error()
            print(f"[Notion] Connected to workspace: {self.workspace_name}")
            return True

        except Exception as e:
            self._set_error(f"Connection failed: {str(e)}")
            return False

    def disconnect(self) -> bool:
        """Disconnect from Notion"""
        self.client = None
        self.status = ConnectorStatus.DISCONNECTED
        return True

    def test_connection(self) -> bool:
        """Verify connection is still valid"""
        if not self.client:
            return False
        try:
            self.client.users.me()
            return True
        except Exception:
            return False

    def sync(self, since: Optional[datetime] = None) -> List[Document]:
        """Sync pages from Notion"""
        if not self.client:
            if not self.connect():
                return []

        self.status = ConnectorStatus.SYNCING
        documents = []

        try:
            max_pages = self.config.settings.get("max_pages", 500)
            print(f"[Notion] Starting sync (max_pages={max_pages})")

            has_more = True
            start_cursor = None
            page_count = 0

            while has_more and page_count < max_pages:
                search_params = {
                    "filter": {"property": "object", "value": "page"},
                    "sort": {"direction": "descending", "timestamp": "last_edited_time"},
                    "page_size": 100
                }

                if start_cursor:
                    search_params["start_cursor"] = start_cursor

                results = self.client.search(**search_params)
                pages = results.get("results", [])
                print(f"[Notion] Found {len(pages)} pages in batch")

                for page in pages:
                    # Skip archived unless configured
                    if page.get("archived") and not self.config.settings.get("include_archived"):
                        continue

                    doc = self._page_to_document(page)
                    if doc:
                        documents.append(doc)
                        # Add any files embedded in this page as separate documents
                        if self._child_documents:
                            documents.extend(self._child_documents)
                            print(f"[Notion]   + {len(self._child_documents)} embedded files as separate documents")
                        page_count += 1
                        print(f"[Notion] Processed page {page_count}: {doc.title[:50] if doc.title else 'Untitled'}")

                        if page_count >= max_pages:
                            break

                has_more = results.get("has_more", False)
                start_cursor = results.get("next_cursor")

            self.config.last_sync = datetime.now(timezone.utc)
            self.status = ConnectorStatus.CONNECTED
            print(f"[Notion] Sync complete: {len(documents)} documents")

        except Exception as e:
            self._set_error(f"Sync failed: {str(e)}")
            import traceback
            traceback.print_exc()

        return documents

    def get_document(self, doc_id: str) -> Optional[Document]:
        """Retrieve a specific page by ID"""
        if not self.client:
            self.connect()

        try:
            page_id = doc_id.replace("notion_", "")
            page = self.client.pages.retrieve(page_id)
            return self._page_to_document(page)
        except Exception as e:
            self._set_error(f"Failed to get document: {str(e)}")
            return None

    def _page_to_document(self, page: Dict[str, Any]) -> Optional[Document]:
        """Convert Notion page to Document object"""
        try:
            page_id = page["id"]
            properties = page.get("properties", {})

            # Parse timestamp early (needed for child documents created during content fetch)
            created_time = page.get("created_time", "")
            timestamp = None
            if created_time:
                try:
                    timestamp = datetime.fromisoformat(created_time.replace("Z", "+00:00"))
                except:
                    pass

            # Reset child document accumulator for this page
            self._child_documents = []
            self._current_page_id = page_id
            self._current_page_url = page.get("url", "")
            self._current_page_timestamp = timestamp

            title = self._extract_title(properties)
            properties_text = self._extract_properties(properties)
            body_content = self._get_page_content(page_id)

            # Combine properties + body into full content
            content_parts = []
            if properties_text:
                content_parts.append(properties_text)
            if body_content.strip():
                content_parts.append(body_content)
            content = "\n\n".join(content_parts)

            if not content.strip() and not title:
                return None

            # Build metadata
            metadata = {
                "page_id": page_id,
                "url": page.get("url", ""),
                "archived": page.get("archived", False),
                "workspace": self.workspace_name
            }

            return Document(
                doc_id=f"notion_{page_id}",
                source="notion",
                content=content,
                title=title or "Untitled",
                metadata=metadata,
                timestamp=timestamp,
                author=self.workspace_name,
                url=page.get("url"),
                doc_type="document"
            )

        except Exception as e:
            print(f"[Notion] Error converting page: {e}")
            return None

    def _extract_title(self, properties: Dict) -> str:
        """Extract title from page properties"""
        for prop_name in ["title", "Title", "Name", "name"]:
            if prop_name in properties:
                prop = properties[prop_name]
                if prop.get("type") == "title":
                    title_array = prop.get("title", [])
                    return "".join([t.get("plain_text", "") for t in title_array])

        for prop in properties.values():
            if prop.get("type") == "title":
                title_array = prop.get("title", [])
                return "".join([t.get("plain_text", "") for t in title_array])

        return ""

    def _extract_properties(self, properties: Dict) -> str:
        """Extract all page properties (beyond title) as structured text"""
        prop_lines = []

        for prop_name, prop in properties.items():
            prop_type = prop.get("type", "")

            # Skip title — already used as the document title
            if prop_type == "title":
                continue

            value = self._extract_property_value(prop_type, prop)
            if value:
                prop_lines.append(f"{prop_name}: {value}")

        return "\n".join(prop_lines)

    def _extract_property_value(self, prop_type: str, prop: Dict) -> str:
        """Extract a single property value based on its type"""
        data = prop.get(prop_type)
        if data is None:
            return ""

        if prop_type == "rich_text":
            return "".join([t.get("plain_text", "") for t in data])
        elif prop_type == "number":
            return str(data) if data is not None else ""
        elif prop_type == "select":
            return data.get("name", "") if data else ""
        elif prop_type == "multi_select":
            return ", ".join([s.get("name", "") for s in data]) if data else ""
        elif prop_type == "status":
            return data.get("name", "") if data else ""
        elif prop_type == "date":
            if not data:
                return ""
            start = data.get("start", "")
            end = data.get("end", "")
            return f"{start} to {end}" if end else start
        elif prop_type == "checkbox":
            return "Yes" if data else "No"
        elif prop_type == "url":
            return data if data else ""
        elif prop_type == "email":
            return data if data else ""
        elif prop_type == "phone_number":
            return data if data else ""
        elif prop_type == "people":
            return ", ".join([p.get("name", p.get("id", "")) for p in data]) if data else ""
        elif prop_type == "relation":
            return ", ".join([r.get("id", "") for r in data]) if data else ""
        elif prop_type == "formula":
            if not data:
                return ""
            f_type = data.get("type", "")
            return str(data.get(f_type, ""))
        elif prop_type == "rollup":
            if not data:
                return ""
            r_type = data.get("type", "")
            return str(data.get(r_type, ""))
        elif prop_type == "files":
            parts = []
            for f in (data or []):
                name = f.get("name", "")
                f_type = f.get("type", "")
                url = ""
                if f_type == "file":
                    url = f.get("file", {}).get("url", "")
                elif f_type == "external":
                    url = f.get("external", {}).get("url", "")
                if url and name:
                    try:
                        ext = os.path.splitext(name)[1].lower()
                        if ext in ('.mp4', '.mov', '.wav', '.mp3', '.m4a', '.webm'):
                            parsed = self._download_and_transcribe_media(url, name)
                        else:
                            parsed = self._download_and_parse_file(url, name)
                        if parsed:
                            self._create_child_document(name, parsed, url)
                            parts.append(f"--- {name} ---\n{parsed}\n--- End of {name} ---")
                        else:
                            parts.append(name)
                    except Exception as e:
                        print(f"[Notion] Error parsing property file {name}: {e}")
                        parts.append(name)
                elif name:
                    parts.append(name)
                elif url:
                    parts.append(url)
            return "\n\n".join(parts) if parts else ""
        elif prop_type in ("created_time", "last_edited_time"):
            return data if data else ""
        elif prop_type == "created_by":
            return data.get("name", data.get("id", "")) if data else ""
        elif prop_type == "last_edited_by":
            return data.get("name", data.get("id", "")) if data else ""
        elif prop_type == "unique_id":
            if not data:
                return ""
            prefix = data.get("prefix", "")
            number = data.get("number", "")
            return f"{prefix}-{number}" if prefix else str(number)

        return ""

    def _get_page_content(self, page_id: str, depth: int = 0) -> str:
        """Fetch and concatenate all blocks from a page"""
        content_parts = []
        max_blocks = self.config.settings.get("max_blocks_per_page", 1000)
        max_depth = self.config.settings.get("max_depth", 5)

        if depth > max_depth:
            return ""

        try:
            has_more = True
            start_cursor = None
            block_count = 0
            first_request = True

            while has_more and block_count < max_blocks:
                params = {"block_id": page_id, "page_size": 100}
                if start_cursor:
                    params["start_cursor"] = start_cursor

                response = self.client.blocks.children.list(**params)
                results = response.get("results", [])

                # Log empty results on first request (helps debug permission issues)
                if first_request and not results:
                    print(f"[Notion] WARNING: blocks.children.list returned 0 blocks for page {page_id}. "
                          f"This may indicate the integration lacks 'Read content' capability. "
                          f"Check integration settings at https://www.notion.so/my-integrations")
                first_request = False

                for block in results:
                    block_type = block.get("type", "unknown")
                    text = self._extract_block_text(block)
                    if text:
                        content_parts.append(text)
                    block_count += 1

                    if block.get("has_children") and block_count < max_blocks:
                        child_content = self._get_page_content(block["id"], depth + 1)
                        if child_content:
                            content_parts.append(child_content)

                    # Handle link_to_page — fetch the linked page's content
                    if block_type == "link_to_page" and depth < max_depth:
                        link_data = block.get("link_to_page", {})
                        target_type = link_data.get("type", "")
                        target_id = link_data.get(target_type, "")
                        if target_id:
                            try:
                                linked_content = self._get_page_content(target_id, depth + 1)
                                if linked_content:
                                    content_parts.append(linked_content)
                            except Exception as e:
                                print(f"[Notion] Error fetching linked page {target_id}: {e}")

                has_more = response.get("has_more", False)
                start_cursor = response.get("next_cursor")

            if depth == 0:
                print(f"[Notion] Page {page_id}: extracted {block_count} blocks, {len(content_parts)} text parts")

        except Exception as e:
            error_str = str(e)
            if "403" in error_str or "restricted" in error_str.lower():
                print(f"[Notion] PERMISSION ERROR for page {page_id}: {e}")
                print(f"[Notion] The integration needs 'Read content' capability. "
                      f"Enable it at https://www.notion.so/my-integrations")
            else:
                print(f"[Notion] Error fetching blocks for {page_id}: {e}")

        return "\n".join(content_parts)

    def _get_rich_text(self, data: Dict) -> str:
        """Extract rich text with hyperlinks preserved"""
        texts = data.get("rich_text", [])
        parts = []
        for t in texts:
            plain = t.get("plain_text", "")
            href = t.get("href")
            if href and plain:
                parts.append(f"{plain} ({href})")
            else:
                parts.append(plain)
        return "".join(parts)

    def _get_media_info(self, block_data: Dict) -> str:
        """Extract URL and caption from media blocks (image, video, audio, file, pdf)"""
        caption = ""
        caption_items = block_data.get("caption", [])
        if caption_items:
            caption = "".join([t.get("plain_text", "") for t in caption_items])

        url = ""
        media_type = block_data.get("type", "")
        if media_type == "external":
            url = block_data.get("external", {}).get("url", "")
        elif media_type == "file":
            url = block_data.get("file", {}).get("url", "")
        elif media_type == "url":
            url = block_data.get("url", "")

        return caption if caption else url

    def _extract_table_row(self, block_data: Dict) -> str:
        """Extract table row cells as pipe-delimited markdown"""
        cells = block_data.get("cells", [])
        cell_texts = []
        for cell in cells:
            cell_text = "".join([t.get("plain_text", "") for t in cell])
            cell_texts.append(cell_text)
        return "| " + " | ".join(cell_texts) + " |"

    def _get_file_url(self, block_data: Dict) -> str:
        """Extract the download URL from a file/pdf block's data"""
        media_type = block_data.get("type", "")
        if media_type == "external":
            return block_data.get("external", {}).get("url", "")
        elif media_type == "file":
            return block_data.get("file", {}).get("url", "")
        return ""

    def _filename_from_url(self, url: str, fallback: str = "file") -> str:
        """Extract filename from URL path, with fallback"""
        if not url:
            return fallback
        try:
            from urllib.parse import urlparse, unquote
            path = urlparse(url).path
            if path:
                name = unquote(path.split("/")[-1])
                # Only use if it looks like a filename with extension
                if name and "." in name:
                    return name
        except Exception:
            pass
        return fallback

    def _get_file_doc_type(self, filename: str) -> str:
        """Return doc_type string based on file extension"""
        ext = os.path.splitext(filename)[1].lower() if filename else ""
        if ext == ".pdf":
            return "pdf"
        elif ext in (".doc", ".docx"):
            return "document"
        elif ext in (".xls", ".xlsx"):
            return "spreadsheet"
        elif ext in (".ppt", ".pptx"):
            return "presentation"
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"):
            return "image"
        elif ext in (".mp4", ".mov", ".webm"):
            return "video"
        elif ext in (".mp3", ".wav", ".m4a"):
            return "audio"
        return "document"

    def _create_child_document(self, filename: str, content: str, source_url: str = "") -> None:
        """Create a child Document for an embedded file and add it to _child_documents"""
        import hashlib
        file_hash = hashlib.md5(f"{self._current_page_id}_{filename}_{source_url}".encode()).hexdigest()[:12]
        doc = Document(
            doc_id=f"notion_file_{self._current_page_id}_{file_hash}",
            source="notion",
            content=content,
            title=filename,
            metadata={
                "parent_page_id": self._current_page_id,
                "parent_page_url": self._current_page_url,
                "file_url": source_url,
                "embedded_in_notion": True
            },
            timestamp=self._current_page_timestamp,
            author=self.workspace_name,
            url=None,
            doc_type=self._get_file_doc_type(filename)
        )
        self._child_documents.append(doc)
        print(f"[Notion] Created child document: {filename} (type={doc.doc_type})")

    def _download_and_transcribe_media(self, url: str, filename: str) -> Optional[str]:
        """Download video/audio from Notion and transcribe with Azure Whisper"""
        if not url:
            return None
        try:
            print(f"[Notion] Downloading media for transcription: {filename}")
            resp = requests.get(url, timeout=120)
            resp.raise_for_status()
            media_bytes = resp.content
            print(f"[Notion] Downloaded {len(media_bytes)} bytes for {filename}")

            from services.knowledge_service import KnowledgeService
            from database.models import SessionLocal
            db = SessionLocal()
            try:
                ks = KnowledgeService(db)
                result = ks.transcribe_audio(media_bytes, filename)
                if result and result.text:
                    print(f"[Notion] Transcribed {len(result.text)} chars from {filename}")
                    return result.text
            finally:
                db.close()
            return None
        except requests.exceptions.RequestException as e:
            print(f"[Notion] Failed to download media {filename}: {e}")
        except Exception as e:
            print(f"[Notion] Error transcribing {filename}: {e}")
        return None

    def _download_and_parse_file(self, url: str, filename: str) -> Optional[str]:
        """Download a file from Notion's S3 URL and parse its content

        Uses Azure Mistral Document AI (mistral-document-ai-2505) as the primary parser,
        which handles PDF, DOCX, PPTX, XLSX, HTML, XML, PNG, JPG and more.
        Falls back to local parsers (PyPDF2, python-docx, etc.) if Mistral is unavailable.
        """
        if not url:
            return None

        ext = os.path.splitext(filename)[1].lower() if filename else ""

        try:
            print(f"[Notion] Downloading file: {filename}")
            resp = requests.get(url, timeout=60)
            resp.raise_for_status()
            file_bytes = resp.content
            print(f"[Notion] Downloaded {len(file_bytes)} bytes for {filename}")

            # Plain text files - decode directly
            if ext in (".txt", ".csv", ".md", ".json"):
                try:
                    return file_bytes.decode("utf-8")
                except UnicodeDecodeError:
                    return file_bytes.decode("latin-1")

            # Try Azure Mistral Document AI first (handles all formats including images)
            parsed = self._parse_with_mistral(file_bytes, filename, ext)
            if parsed:
                return parsed

            # Fallback to local parsers for common Office formats
            print(f"[Notion] Mistral unavailable, trying local parser for {filename}")
            return self._parse_with_local(file_bytes, filename, ext)

        except requests.exceptions.RequestException as e:
            print(f"[Notion] Failed to download {filename}: {e}")
        except Exception as e:
            print(f"[Notion] Error parsing {filename}: {e}")

        return None

    def _parse_with_mistral(self, file_bytes: bytes, filename: str, ext: str) -> Optional[str]:
        """Parse file using Azure Mistral Document AI (mistral-document-ai-2505)"""
        try:
            from parsers.azure_doc_parser import AzureDocumentParser
            parser = AzureDocumentParser()

            # Mistral parser expects a file path, so write to temp file
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            try:
                result = parser.parse(tmp_path)
                if result and result.get("content"):
                    content = result["content"].replace("\x00", "")
                    print(f"[Notion] Parsed {filename} with Mistral Document AI: {len(content)} chars")
                    return content
            finally:
                os.unlink(tmp_path)

        except Exception as e:
            print(f"[Notion] Mistral Document AI not available for {filename}: {e}")

        return None

    def _parse_with_local(self, file_bytes: bytes, filename: str, ext: str) -> Optional[str]:
        """Fallback: parse file using local libraries (PyPDF2, python-docx, etc.)"""
        try:
            if ext == ".pdf":
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
                parts = [p.extract_text() for p in pdf_reader.pages if p.extract_text()]
                result = "\n\n".join(parts)
                print(f"[Notion] Parsed PDF {filename} locally: {len(result)} chars, {len(pdf_reader.pages)} pages")
                return result

            elif ext in (".docx", ".doc"):
                from docx import Document as DocxDocument
                doc = DocxDocument(io.BytesIO(file_bytes))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                result = "\n\n".join(parts)
                print(f"[Notion] Parsed Word {filename} locally: {len(result)} chars")
                return result

            elif ext in (".pptx", ".ppt"):
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = f"\n--- Slide {i} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += shape.text + "\n"
                    parts.append(slide_text)
                result = "\n".join(parts)
                print(f"[Notion] Parsed PowerPoint {filename} locally: {len(result)} chars, {len(prs.slides)} slides")
                return result

            elif ext in (".xlsx", ".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    parts.append(f"\n--- Sheet: {sheet_name} ---\n")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join(str(c) if c is not None else "" for c in row)
                        if row_text.strip():
                            parts.append(row_text)
                result = "\n".join(parts)
                print(f"[Notion] Parsed Excel {filename} locally: {len(result)} chars, {len(wb.sheetnames)} sheets")
                return result

            elif ext in (".html", ".htm"):
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(file_bytes, "html.parser")
                result = soup.get_text(separator="\n", strip=True)
                print(f"[Notion] Parsed HTML {filename} locally: {len(result)} chars")
                return result

            else:
                print(f"[Notion] No local parser for {filename} ({ext})")

        except Exception as e:
            print(f"[Notion] Local parse error for {filename}: {e}")

        return None

    def _extract_block_text(self, block: Dict) -> str:
        """Extract text from any Notion block type"""
        block_type = block.get("type", "")
        block_data = block.get(block_type, {})

        # --- Text blocks (rich_text based) ---
        if block_type == "paragraph":
            return self._get_rich_text(block_data)
        elif block_type == "heading_1":
            return f"# {self._get_rich_text(block_data)}"
        elif block_type == "heading_2":
            return f"## {self._get_rich_text(block_data)}"
        elif block_type == "heading_3":
            return f"### {self._get_rich_text(block_data)}"
        elif block_type == "bulleted_list_item":
            return f"- {self._get_rich_text(block_data)}"
        elif block_type == "numbered_list_item":
            return f"1. {self._get_rich_text(block_data)}"
        elif block_type == "to_do":
            checked = "x" if block_data.get("checked") else " "
            return f"- [{checked}] {self._get_rich_text(block_data)}"
        elif block_type == "code":
            language = block_data.get("language", "text")
            code = self._get_rich_text(block_data)
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            result = f"```{language}\n{code}\n```"
            if caption:
                result += f"\n{caption}"
            return result
        elif block_type == "quote":
            return f"> {self._get_rich_text(block_data)}"
        elif block_type == "divider":
            return "---"
        elif block_type == "toggle":
            return f"▸ {self._get_rich_text(block_data)}"
        elif block_type == "callout":
            icon = ""
            icon_data = block_data.get("icon")
            if icon_data:
                if icon_data.get("type") == "emoji":
                    icon = icon_data.get("emoji", "")
            text = self._get_rich_text(block_data)
            return f"{icon} {text}".strip() if icon else text

        # --- Table blocks ---
        elif block_type == "table":
            # Table block itself has no text; rows come via child block recursion
            return ""
        elif block_type == "table_row":
            return self._extract_table_row(block_data)

        # --- Media blocks (download + parse/transcribe) ---
        elif block_type == "image":
            url = self._get_file_url(block_data)
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            img_name = self._filename_from_url(url, "image.png")
            if not img_name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
                img_name = img_name + ".png"
            try:
                parsed = self._download_and_parse_file(url, img_name) if url else None
                if parsed:
                    self._create_child_document(img_name, parsed, url)
                    header = f"[Image: {caption}]" if caption else "[Image]"
                    return f"{header}\n{parsed}"
            except Exception as e:
                print(f"[Notion] Image OCR failed: {e}")
            return f"[Image: {caption or url}]" if (caption or url) else "[Image]"

        elif block_type == "video":
            url = self._get_file_url(block_data)
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            vid_name = self._filename_from_url(url, "video.mp4")
            try:
                parsed = self._download_and_transcribe_media(url, vid_name) if url else None
                if parsed:
                    self._create_child_document(vid_name, parsed, url)
                    header = f"[Video: {caption}]" if caption else "[Video]"
                    return f"{header}\n--- Video Transcript ---\n{parsed}\n--- End Transcript ---"
            except Exception as e:
                print(f"[Notion] Video transcription failed: {e}")
            return f"[Video: {caption or url}]" if (caption or url) else "[Video]"

        elif block_type == "audio":
            url = self._get_file_url(block_data)
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            aud_name = self._filename_from_url(url, "audio.mp3")
            try:
                parsed = self._download_and_transcribe_media(url, aud_name) if url else None
                if parsed:
                    self._create_child_document(aud_name, parsed, url)
                    header = f"[Audio: {caption}]" if caption else "[Audio]"
                    return f"{header}\n--- Audio Transcript ---\n{parsed}\n--- End Transcript ---"
            except Exception as e:
                print(f"[Notion] Audio transcription failed: {e}")
            return f"[Audio: {caption or url}]" if (caption or url) else "[Audio]"

        elif block_type == "file":
            url = self._get_file_url(block_data)
            name = block_data.get("name", "")
            if not name and url:
                name = self._filename_from_url(url, "file")
            ext = os.path.splitext(name)[1].lower() if name else ""
            try:
                if ext in ('.mp4', '.mov', '.wav', '.mp3', '.m4a', '.webm'):
                    parsed = self._download_and_transcribe_media(url, name) if url else None
                else:
                    parsed = self._download_and_parse_file(url, name) if url else None
                if parsed:
                    self._create_child_document(name, parsed, url)
                    return f"--- Attached File: {name} ---\n{parsed}\n--- End of {name} ---"
            except Exception as e:
                print(f"[Notion] File parse failed for {name}: {e}")
            return f"[File: {name or url}]" if (name or url) else "[File]"
        elif block_type == "pdf":
            url = self._get_file_url(block_data)
            caption = self._get_media_info(block_data)
            # Derive filename from caption or URL
            pdf_name = caption if caption and caption.endswith(".pdf") else ""
            if not pdf_name and url:
                # Extract filename from S3 URL path
                try:
                    from urllib.parse import urlparse, unquote
                    path = urlparse(url).path
                    pdf_name = unquote(path.split("/")[-1]) if path else "document.pdf"
                except:
                    pdf_name = "document.pdf"
            parsed = self._download_and_parse_file(url, pdf_name) if url else None
            if parsed:
                self._create_child_document(pdf_name, parsed, url)
                return f"--- Attached PDF: {pdf_name} ---\n{parsed}\n--- End of {pdf_name} ---"
            return f"[PDF: {caption or url}]" if (caption or url) else "[PDF]"

        # --- Link/reference blocks ---
        elif block_type == "bookmark":
            url = block_data.get("url", "")
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            if caption:
                return f"[Bookmark: {url}] {caption}"
            return f"[Bookmark: {url}]" if url else ""
        elif block_type == "embed":
            url = block_data.get("url", "")
            caption = "".join([t.get("plain_text", "") for t in block_data.get("caption", [])])
            if caption:
                return f"[Embed: {url}] {caption}"
            return f"[Embed: {url}]" if url else ""
        elif block_type == "link_preview":
            url = block_data.get("url", "")
            return f"[Link: {url}]" if url else ""
        elif block_type == "link_to_page":
            page_type = block_data.get("type", "")
            target_id = block_data.get(page_type, "")
            return f"\n## Linked Page\n"
        elif block_type == "child_page":
            child_title = block_data.get("title", "")
            return f"\n## Child Page: {child_title}\n"
        elif block_type == "child_database":
            db_title = block_data.get("title", "")
            return f"[Database: {db_title}]"

        # --- Structural blocks (content comes from children via recursion) ---
        elif block_type in ("synced_block", "column_list", "column"):
            return ""

        # --- Other blocks ---
        elif block_type == "equation":
            expression = block_data.get("expression", "")
            return f"$$ {expression} $$" if expression else ""
        elif block_type == "template":
            return f"[Template: {self._get_rich_text(block_data)}]"
        elif block_type == "table_of_contents":
            return "[Table of Contents]"
        elif block_type == "breadcrumb":
            return ""

        # Fallback: try to extract rich_text if present
        if "rich_text" in block_data:
            return self._get_rich_text(block_data)

        return ""

    @classmethod
    def get_auth_url(cls, redirect_uri: str, state: str) -> str:
        """Generate Notion OAuth authorization URL"""
        client_id = os.getenv("NOTION_CLIENT_ID", "")
        return (
            f"https://api.notion.com/v1/oauth/authorize"
            f"?client_id={client_id}"
            f"&redirect_uri={redirect_uri}"
            f"&response_type=code"
            f"&owner=user"
            f"&state={state}"
        )

    @classmethod
    def exchange_code(cls, code: str, redirect_uri: str) -> Dict[str, Any]:
        """Exchange authorization code for access token"""
        client_id = os.getenv("NOTION_CLIENT_ID", "")
        client_secret = os.getenv("NOTION_CLIENT_SECRET", "")

        credentials = base64.b64encode(f"{client_id}:{client_secret}".encode()).decode()

        response = requests.post(
            "https://api.notion.com/v1/oauth/token",
            headers={
                "Authorization": f"Basic {credentials}",
                "Content-Type": "application/json"
            },
            json={
                "grant_type": "authorization_code",
                "code": code,
                "redirect_uri": redirect_uri
            }
        )

        data = response.json()

        if "error" in data:
            raise Exception(f"OAuth failed: {data.get('error_description', data['error'])}")

        return {
            "access_token": data.get("access_token"),
            "workspace_id": data.get("workspace_id"),
            "workspace_name": data.get("workspace_name"),
            "bot_id": data.get("bot_id")
        }
