"""
OneDrive Connector
Connects to Microsoft OneDrive/SharePoint to extract PowerPoint, Excel, and Word files.
"""

import os
import io
import mimetypes
from datetime import datetime
from typing import List, Dict, Optional, Any

from .base_connector import BaseConnector, ConnectorConfig, ConnectorStatus, Document

# Note: Requires msal and requests
# pip install msal requests

try:
    import msal
    import requests
    ONEDRIVE_AVAILABLE = True
except ImportError:
    ONEDRIVE_AVAILABLE = False


class OneDriveConnector(BaseConnector):
    """
    OneDrive connector for extracting Microsoft Office files.

    Extracts:
    - PowerPoint presentations (.pptx, .ppt)
    - Excel spreadsheets (.xlsx, .xls)
    - Word documents (.docx, .doc)
    - PDFs
    """

    CONNECTOR_TYPE = "onedrive"
    REQUIRED_CREDENTIALS = ["access_token", "refresh_token"]
    OPTIONAL_SETTINGS = {
        "folder_ids": [],  # Specific folders to sync (empty = root)
        "file_types": [
            ".pptx", ".ppt", ".xlsx", ".xls", ".docx", ".doc", ".pdf",
            ".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml",
            ".rtf", ".odt", ".ods", ".odp", ".tex", ".log", ".yaml", ".yml"
        ],
        "max_file_size_mb": 50,
        "include_shared": True
    }

    # Microsoft Graph API endpoint
    GRAPH_ENDPOINT = "https://graph.microsoft.com/v1.0"

    def __init__(self, config: ConnectorConfig):
        super().__init__(config)
        self.access_token = None

    # Scopes for personal Microsoft accounts (OneDrive Personal)
    # Note: MSAL automatically adds offline_access, openid, profile
    SCOPES = ["Files.ReadWrite", "User.Read"]

    @classmethod
    def _get_oauth_config(cls) -> Dict:
        return {
            "client_id": os.getenv("MICROSOFT_CLIENT_ID", ""),
            "client_secret": os.getenv("MICROSOFT_CLIENT_SECRET", ""),
            "tenant": os.getenv("MICROSOFT_TENANT_ID", "consumers"),
            "redirect_uri": os.getenv("MICROSOFT_REDIRECT_URI", "http://localhost:5003/api/integrations/onedrive/callback")
        }

    @classmethod
    def get_auth_url(cls, redirect_uri: str, state: str) -> str:
        """Generate Microsoft OAuth authorization URL"""
        if not ONEDRIVE_AVAILABLE:
            raise ImportError("MSAL not installed")

        config = cls._get_oauth_config()

        # Create MSAL app (ConfidentialClient for server-side apps with client_secret)
        app = msal.ConfidentialClientApplication(
            config["client_id"],
            authority=f"https://login.microsoftonline.com/{config['tenant']}",
            client_credential=config["client_secret"]
        )

        auth_url = app.get_authorization_request_url(
            cls.SCOPES,
            state=state,
            redirect_uri=redirect_uri
        )

        return auth_url

    @classmethod
    def exchange_code_for_tokens(cls, code: str, redirect_uri: str):
        """Exchange authorization code for access token"""
        try:
            if not ONEDRIVE_AVAILABLE:
                return None, "MSAL not installed. Run: pip install msal"

            config = cls._get_oauth_config()

            # Create MSAL app (ConfidentialClient for server-side apps with client_secret)
            app = msal.ConfidentialClientApplication(
                config["client_id"],
                authority=f"https://login.microsoftonline.com/{config['tenant']}",
                client_credential=config["client_secret"]
            )

            # Exchange code for token
            result = app.acquire_token_by_authorization_code(
                code,
                scopes=cls.SCOPES,
                redirect_uri=redirect_uri
            )

            if "error" in result:
                return None, result.get("error_description", result["error"])

            tokens = {
                "access_token": result["access_token"],
                "refresh_token": result.get("refresh_token"),
                "expires_in": result.get("expires_in")
            }

            return tokens, None

        except Exception as e:
            return None, str(e)

    def connect(self) -> bool:
        """Connect to OneDrive"""
        if not ONEDRIVE_AVAILABLE:
            self._set_error("MSAL not installed. Run: pip install msal")
            return False

        try:
            self.status = ConnectorStatus.CONNECTING
            self.access_token = self.config.credentials.get("access_token")

            # Test connection by getting user profile
            response = requests.get(
                f"{self.GRAPH_ENDPOINT}/me",
                headers={"Authorization": f"Bearer {self.access_token}"}
            )

            if response.status_code != 200:
                self._set_error(f"Connection failed: {response.text}")
                return False

            user_data = response.json()
            self.sync_stats["user"] = user_data.get("displayName", "Unknown")
            print(f"[OneDrive] Connected as: {user_data.get('displayName')} ({user_data.get('userPrincipalName', user_data.get('mail', 'unknown'))})", flush=True)

            # Check drive info
            drive_response = requests.get(
                f"{self.GRAPH_ENDPOINT}/me/drive",
                headers={"Authorization": f"Bearer {self.access_token}"}
            )
            if drive_response.status_code == 200:
                drive_data = drive_response.json()
                drive_type = drive_data.get("driveType", "unknown")
                drive_id = drive_data.get("id", "unknown")
                quota = drive_data.get("quota", {})
                used = quota.get("used", 0) / (1024*1024)
                total = quota.get("total", 0) / (1024*1024*1024)
                print(f"[OneDrive] Drive type: {drive_type}, id: {drive_id}, used: {used:.1f}MB / {total:.1f}GB", flush=True)
                self._drive_id = drive_id
            else:
                print(f"[OneDrive] WARNING: /me/drive returned {drive_response.status_code}: {drive_response.text[:200]}", flush=True)
                self._drive_id = None

            self.status = ConnectorStatus.CONNECTED
            self._clear_error()
            return True

        except Exception as e:
            self._set_error(f"Failed to connect: {str(e)}")
            return False

    def disconnect(self) -> bool:
        """Disconnect from OneDrive"""
        self.access_token = None
        self.status = ConnectorStatus.DISCONNECTED
        return True

    def test_connection(self) -> bool:
        """Test OneDrive connection"""
        if not self.access_token:
            return False

        try:
            response = requests.get(
                f"{self.GRAPH_ENDPOINT}/me/drive",
                headers={"Authorization": f"Bearer {self.access_token}"}
            )
            return response.status_code == 200
        except Exception:
            return False

    def sync(self, since: Optional[datetime] = None) -> List[Document]:
        """Sync files from OneDrive"""
        if not self.access_token:
            self.connect()

        if self.status != ConnectorStatus.CONNECTED:
            return []

        self.status = ConnectorStatus.SYNCING
        documents = []

        try:
            # Get folders to sync
            folder_ids = self.config.settings.get("folder_ids", [])

            if folder_ids:
                # Sync specific folders
                for folder_id in folder_ids:
                    folder_docs = self._sync_folder(folder_id, since)
                    documents.extend(folder_docs)
            else:
                # Sync root folder
                folder_docs = self._sync_folder("root", since)
                documents.extend(folder_docs)

            # Update stats
            self.sync_stats = {
                "documents_synced": len(documents),
                "sync_time": datetime.now().isoformat()
            }

            self.config.last_sync = datetime.now()
            self.status = ConnectorStatus.CONNECTED

        except Exception as e:
            self._set_error(f"Sync failed: {str(e)}")
            raise  # Propagate to caller so sync reports as error

        return documents

    def _sync_folder(self, folder_id: str, since: Optional[datetime]) -> List[Document]:
        """Sync files from a folder recursively"""
        documents = []

        try:
            # Build URL - try multiple endpoint formats for compatibility
            urls_to_try = []
            if folder_id == "root":
                urls_to_try.append(f"{self.GRAPH_ENDPOINT}/me/drive/root/children")
                # Fallback: use drive ID if available
                if hasattr(self, '_drive_id') and self._drive_id:
                    urls_to_try.append(f"{self.GRAPH_ENDPOINT}/drives/{self._drive_id}/root/children")
                urls_to_try.append(f"{self.GRAPH_ENDPOINT}/me/drive/items/root/children")
            else:
                urls_to_try.append(f"{self.GRAPH_ENDPOINT}/me/drive/items/{folder_id}/children")

            # Try each URL format until one works
            url = None
            response = None
            for try_url in urls_to_try:
                print(f"[OneDrive] Trying: {try_url}", flush=True)
                response = requests.get(
                    try_url,
                    headers={"Authorization": f"Bearer {self.access_token}"}
                )
                if response.status_code == 200:
                    url = try_url
                    print(f"[OneDrive] Success with: {try_url}", flush=True)
                    break
                else:
                    print(f"[OneDrive] {try_url} returned {response.status_code}: {response.text[:200]}", flush=True)

            if not url or response.status_code != 200:
                error_text = response.text if response else "No response"
                try:
                    error_data = response.json().get("error", {})
                    error_msg = error_data.get("message", error_text[:200])
                except Exception:
                    error_msg = f"HTTP {response.status_code}: {error_text[:200]}"
                raise Exception(f"OneDrive API error: {error_msg}")

            # Process first page and then paginate
            while True:
                data = response.json()
                items = data.get("value", [])
                print(f"[OneDrive] Folder {folder_id}: got {len(items)} items", flush=True)

                for item in items:
                    if "folder" in item:
                        subfolder_name = item.get("name", "unknown")
                        print(f"[OneDrive] Entering subfolder: {subfolder_name}", flush=True)
                        subfolder_docs = self._sync_folder(item["id"], since)
                        documents.extend(subfolder_docs)

                    elif "file" in item:
                        name = item.get("name", "")
                        file_types = self.config.settings.get("file_types", self.OPTIONAL_SETTINGS["file_types"])

                        if any(name.lower().endswith(ext) for ext in file_types):
                            print(f"[OneDrive] Found matching file: {name}", flush=True)
                            size_mb = item.get("size", 0) / (1024 * 1024)
                            max_size = self.config.settings.get("max_file_size_mb", 50)

                            if size_mb > max_size:
                                print(f"[OneDrive] Skipping {name} - too large ({size_mb:.1f}MB)")
                                continue

                            modified = datetime.fromisoformat(item["lastModifiedDateTime"].replace("Z", "+00:00"))

                            if since and modified < since:
                                continue

                            doc = self._download_and_parse(item)
                            if doc:
                                documents.append(doc)

                # Check for next page
                next_link = data.get("@odata.nextLink")
                if not next_link:
                    break
                print(f"[OneDrive] Fetching next page...", flush=True)
                response = requests.get(
                    next_link,
                    headers={"Authorization": f"Bearer {self.access_token}"}
                )
                if response.status_code != 200:
                    print(f"[OneDrive] Pagination failed: {response.status_code}", flush=True)
                    break

        except Exception as e:
            print(f"[OneDrive] Error syncing folder {folder_id}: {e}", flush=True)

        return documents

    def _download_and_parse(self, item: Dict) -> Optional[Document]:
        """Download and parse a file"""
        try:
            file_id = item["id"]
            name = item.get("name", "unknown")
            download_url = item.get("@microsoft.graph.downloadUrl")

            if not download_url:
                print(f"[OneDrive] No download URL for {name}")
                return None

            print(f"[OneDrive] Downloading {name}...")

            # Download file
            response = requests.get(download_url)

            if response.status_code != 200:
                print(f"[OneDrive] Failed to download {name}")
                return None

            file_content = response.content

            # Parse based on file type
            content_text = self._parse_file(name, file_content)

            if not content_text:
                return None

            # Create document
            modified = datetime.fromisoformat(item["lastModifiedDateTime"].replace("Z", "+00:00"))

            return Document(
                doc_id=f"onedrive_{file_id}",
                source="onedrive",
                content=content_text,
                title=name,
                metadata={
                    "file_id": file_id,
                    "size": item.get("size", 0),
                    "path": item.get("parentReference", {}).get("path", ""),
                    "web_url": item.get("webUrl")
                },
                timestamp=modified,
                author=item.get("createdBy", {}).get("user", {}).get("displayName"),
                url=item.get("webUrl"),
                doc_type=self._get_doc_type(name)
            )

        except Exception as e:
            print(f"[OneDrive] Error parsing {item.get('name')}: {e}")
            return None

    def _parse_file(self, filename: str, content: bytes) -> Optional[str]:
        """Parse file content based on type. Tries Mistral Document AI first, falls back to local parsers."""
        try:
            lower_name = filename.lower()

            # Plain text files - decode directly (no need for Mistral)
            if lower_name.endswith((
                ".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml",
                ".rtf", ".log", ".yaml", ".yml", ".tex"
            )):
                return self._parse_text(content)

            # For binary document files, try Mistral Document AI first
            if lower_name.endswith((".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".odt", ".ods", ".odp")):
                try:
                    from parsers.document_parser import DocumentParser
                    parser = DocumentParser()
                    parsed = parser.parse_file_bytes(content, filename)
                    if parsed:
                        print(f"[OneDrive] Parsed {filename} with Mistral Document AI: {len(parsed)} chars")
                        return parsed
                except Exception as e:
                    print(f"[OneDrive] Mistral unavailable for {filename}: {e}")

            # Fallback to local parsers
            if lower_name.endswith((".pptx", ".ppt")):
                return self._parse_powerpoint(content)
            elif lower_name.endswith((".xlsx", ".xls")):
                return self._parse_excel(content)
            elif lower_name.endswith((".docx", ".doc")):
                return self._parse_word(content)
            elif lower_name.endswith(".pdf"):
                return self._parse_pdf(content)

            return None

        except Exception as e:
            print(f"[OneDrive] Parse error for {filename}: {e}")
            return None

    def _parse_text(self, content: bytes) -> Optional[str]:
        """Parse plain text file"""
        try:
            # Try UTF-8 first, then fallback to latin-1
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1")
        except Exception as e:
            print(f"[OneDrive] Text parse error: {e}")
            return None

    def _parse_powerpoint(self, content: bytes) -> Optional[str]:
        """Parse PowerPoint file"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                text_parts.append(slide_text)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] PowerPoint parse error: {e}")
            return None

    def _parse_excel(self, content: bytes) -> Optional[str]:
        """Parse Excel file"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")

                # Get used range
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] Excel parse error: {e}")
            return None

    def _parse_word(self, content: bytes) -> Optional[str]:
        """Parse Word document"""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            return "\n\n".join(paragraphs)

        except Exception as e:
            print(f"[OneDrive] Word parse error: {e}")
            return None

    def _parse_pdf(self, content: bytes) -> Optional[str]:
        """Parse PDF file"""
        try:
            import PyPDF2

            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text_parts = []

            for page in pdf_reader.pages:
                text_parts.append(page.extract_text())

            return "\n\n".join(text_parts)

        except Exception as e:
            print(f"[OneDrive] PDF parse error: {e}")
            return None

    def _get_doc_type(self, filename: str) -> str:
        """Get document type based on filename"""
        lower_name = filename.lower()

        if lower_name.endswith((".pptx", ".ppt", ".odp")):
            return "presentation"
        elif lower_name.endswith((".xlsx", ".xls", ".csv", ".ods")):
            return "spreadsheet"
        elif lower_name.endswith((".docx", ".doc", ".odt", ".rtf")):
            return "document"
        elif lower_name.endswith(".pdf"):
            return "pdf"
        elif lower_name.endswith((".txt", ".md", ".log", ".tex")):
            return "text"
        elif lower_name.endswith((".html", ".htm", ".xml")):
            return "webpage"
        elif lower_name.endswith((".json", ".yaml", ".yml")):
            return "data"

        return "file"

    def get_document(self, doc_id: str) -> Optional[Document]:
        """Get a specific document by ID"""
        # Implementation would fetch single file by ID
        return None
