"""
Azure GPT-4o Document Parser
Uses Azure OpenAI GPT-4o (vision) for intelligent document parsing.
For images: sends directly as base64 image.
For PDFs/Office docs: extracts text via traditional parsers, then uses GPT-4o to structure it.
"""

import os
import base64
from pathlib import Path
from typing import Dict, Optional
import warnings

warnings.filterwarnings('ignore')

try:
    from openai import AzureOpenAI
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False

# Traditional parsers for text extraction before GPT-4o structuring
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


SYSTEM_PROMPT = """You are a document parser. Given the extracted text from a document, produce a clean, well-structured version that preserves all information.

For each document, provide:
1. Main text content (preserve structure and formatting)
2. Tables (extract data in markdown table format)
3. Lists (preserve bullet points and numbering)
4. Headers and sections (maintain hierarchy)
5. Important metadata (dates, names, numbers)

Return clean, well-structured markdown that preserves the document's organization.
Do NOT add commentary — only return the structured document content."""

IMAGE_SYSTEM_PROMPT = """You are a document parser. Extract ALL text and visual content from this image.

Provide:
1. All visible text (preserve structure)
2. Tables (in markdown format)
3. Charts/graphs (describe data and values)
4. Any other relevant content

Return clean, well-structured markdown. Do NOT add commentary."""


class AzureDocumentParser:
    """Parse documents using Azure OpenAI GPT-4o"""

    # Class-level circuit breaker
    _consecutive_failures = 0
    _circuit_open = False
    _FAILURE_THRESHOLD = 3

    def __init__(self, config=None):
        self.config = config
        self.client = None
        self.model = None

        if not HAS_OPENAI:
            raise ImportError("openai not installed. Run: pip install openai")

        self._initialize_client()

        self.supported_formats = [
            '.pdf', '.pptx', '.xlsx', '.docx',
            '.txt', '.html', '.xml',
            '.png', '.jpg', '.jpeg'
        ]

        # Image formats that GPT-4o can process directly via vision
        self.image_formats = ['.png', '.jpg', '.jpeg']

    def _initialize_client(self):
        """Initialize Azure OpenAI client"""
        if self.config:
            api_key = getattr(self.config, 'AZURE_OPENAI_API_KEY', None)
            endpoint = getattr(self.config, 'AZURE_OPENAI_ENDPOINT', None)
            api_version = getattr(self.config, 'AZURE_API_VERSION', None)
            self.model = getattr(self.config, 'AZURE_CHAT_DEPLOYMENT', None)
        else:
            api_key = os.getenv('AZURE_OPENAI_API_KEY')
            endpoint = os.getenv('AZURE_OPENAI_ENDPOINT')
            api_version = os.getenv('AZURE_API_VERSION', '2024-12-01-preview')
            self.model = os.getenv('AZURE_CHAT_DEPLOYMENT', 'gpt-5-chat')

        if not api_key or not endpoint:
            raise ValueError("AZURE_OPENAI_API_KEY and AZURE_OPENAI_ENDPOINT must be set")

        self.client = AzureOpenAI(
            azure_endpoint=endpoint,
            api_key=api_key,
            api_version=api_version
        )
        print(f"✓ Azure GPT-4o document parser initialized (model: {self.model})")

    def can_parse(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in self.supported_formats

    def _extract_raw_text(self, file_path: str, ext: str) -> Optional[str]:
        """Extract raw text from a document using traditional parsers."""
        try:
            if ext == '.pdf' and HAS_PDF:
                parts = []
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text and text.strip():
                            parts.append(text.replace('\x00', '').strip())
                return '\n\n'.join(parts) if parts else None

            elif ext == '.pptx' and HAS_PPTX:
                parts = []
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        parts.append(f"[Slide {i}]\n" + '\n'.join(slide_text))
                return '\n\n'.join(parts) if parts else None

            elif ext == '.docx' and HAS_DOCX:
                doc = DocxDocument(file_path)
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_text = ' | '.join(c.text.strip() for c in row.cells if c.text.strip())
                        if row_text:
                            parts.append(row_text)
                return '\n\n'.join(parts) if parts else None

        except Exception as e:
            print(f"  ⚠ Raw text extraction failed for {Path(file_path).name}: {e}")

        return None

    def parse(self, file_path: str) -> Optional[Dict]:
        """
        Parse a document using GPT-4o.
        - Images: sent directly via vision API
        - PDFs/Office docs: text extracted first, then structured by GPT-4o
        - Text files: read directly (no GPT-4o needed)
        """
        if not os.path.exists(file_path):
            print(f"  ⚠ File not found: {file_path}")
            return None

        ext = Path(file_path).suffix.lower()
        if ext not in self.supported_formats:
            print(f"  ⚠ Unsupported file format: {ext}")
            return None

        if AzureDocumentParser._circuit_open:
            return None

        file_name = Path(file_path).name
        file_size = os.path.getsize(file_path)

        # Skip GPT-4o for large files (>10MB) — too slow and expensive
        if file_size > 10 * 1024 * 1024 and ext not in self.image_formats:
            print(f"  ⚠ {file_name} is {file_size / (1024*1024):.1f}MB — skipping GPT-4o (too large), using traditional parser")
            return None

        try:
            # Text files: read directly, no GPT-4o needed
            if ext in ['.txt', '.html', '.xml']:
                with open(file_path, 'rb') as f:
                    content = f.read().decode('utf-8', errors='ignore')
                return {
                    'content': content,
                    'raw_content': content,
                    'metadata': {
                        'file_type': ext.lstrip('.'),
                        'file_name': file_name,
                        'total_chars': len(content),
                        'parser': 'direct_read',
                        'model': 'none'
                    }
                }

            # Images: send directly to GPT-4o vision
            if ext in self.image_formats:
                return self._parse_image(file_path, file_name, ext)

            # PDFs and Office docs: extract text, then structure with GPT-4o
            return self._parse_document(file_path, file_name, ext)

        except Exception as e:
            AzureDocumentParser._consecutive_failures += 1
            if AzureDocumentParser._consecutive_failures >= AzureDocumentParser._FAILURE_THRESHOLD:
                AzureDocumentParser._circuit_open = True
                print(f"  ✗ GPT-4o parser failed {AzureDocumentParser._consecutive_failures}x — disabling. Error: {e}")
            else:
                print(f"  ✗ Error parsing {file_name} with GPT-4o: {e}")
            return None

    def _parse_image(self, file_path: str, file_name: str, ext: str) -> Optional[Dict]:
        """Parse an image using GPT-4o vision."""
        print(f"  📄 Parsing {file_name} with GPT-4o vision...")

        with open(file_path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode('utf-8')

        mime = 'image/png' if ext == '.png' else 'image/jpeg'

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": IMAGE_SYSTEM_PROMPT},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Extract all content from this image:"},
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"}}
                    ]
                }
            ],
            temperature=0.1,
            max_tokens=16000
        )

        parsed = response.choices[0].message.content
        if not parsed or len(parsed.strip()) < 10:
            print(f"  ⚠ No meaningful content from {file_name}")
            return None

        AzureDocumentParser._consecutive_failures = 0
        return {
            'content': parsed,
            'raw_content': parsed,
            'metadata': {
                'file_type': ext.lstrip('.'),
                'file_name': file_name,
                'total_chars': len(parsed),
                'parser': 'gpt4o_vision',
                'model': self.model,
                'tokens_used': response.usage.total_tokens if hasattr(response, 'usage') else None
            }
        }

    def _parse_document(self, file_path: str, file_name: str, ext: str) -> Optional[Dict]:
        """Parse PDF/Office doc: extract text traditionally, then structure with GPT-4o."""
        print(f"  📄 Parsing {file_name} with GPT-4o (text extraction + structuring)...")

        # Step 1: Extract raw text using traditional parsers
        raw_text = self._extract_raw_text(file_path, ext)

        if not raw_text or len(raw_text.strip()) < 10:
            print(f"  ⚠ Could not extract text from {file_name} for GPT-4o structuring")
            return None

        # Step 2: Send to GPT-4o for intelligent structuring
        # Truncate to avoid token limits (GPT-4o can handle ~128K tokens ≈ 500K chars)
        text_for_gpt = raw_text[:200000]

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {
                    "role": "user",
                    "content": f"Structure and clean the following extracted text from a {ext.lstrip('.')} file named '{file_name}':\n\n{text_for_gpt}"
                }
            ],
            temperature=0.1,
            max_tokens=16000
        )

        parsed = response.choices[0].message.content
        if not parsed or len(parsed.strip()) < 10:
            # Fall back to raw text if GPT-4o returned nothing useful
            print(f"  ⚠ GPT-4o returned insufficient output, using raw text for {file_name}")
            return {
                'content': raw_text,
                'raw_content': raw_text,
                'metadata': {
                    'file_type': ext.lstrip('.'),
                    'file_name': file_name,
                    'total_chars': len(raw_text),
                    'parser': 'traditional_fallback',
                    'model': 'none'
                }
            }

        AzureDocumentParser._consecutive_failures = 0
        return {
            'content': parsed,
            'raw_content': raw_text,
            'metadata': {
                'file_type': ext.lstrip('.'),
                'file_name': file_name,
                'total_chars': len(parsed),
                'raw_chars': len(raw_text),
                'parser': 'gpt4o_structured',
                'model': self.model,
                'tokens_used': response.usage.total_tokens if hasattr(response, 'usage') else None
            }
        }

    def parse_batch(self, file_paths: list) -> Dict[str, Optional[Dict]]:
        results = {}
        for file_path in file_paths:
            results[file_path] = self.parse(file_path)
        return results
