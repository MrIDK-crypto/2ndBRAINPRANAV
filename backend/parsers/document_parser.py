"""
Document Parser for Office Files
Primary: Azure GPT-4o for intelligent document parsing
Fallback: Docling (local, free, no API keys)
Last resort: Traditional parsers (PyPDF2, python-docx, openpyxl, python-pptx)
"""

import os
import csv
import io
import tempfile
from pathlib import Path
from typing import Dict, List, Optional
import warnings
warnings.filterwarnings('ignore')

# Try to import Azure GPT-4o parser (primary)
try:
    from parsers.azure_doc_parser import AzureDocumentParser
    HAS_AZURE_DOC_AI = True
except ImportError:
    HAS_AZURE_DOC_AI = False

# Try to import Docling parser (fallback)
try:
    from parsers.docling_parser import DoclingParser
    HAS_DOCLING = True
except ImportError:
    HAS_DOCLING = False

# PDF parsing (fallback)
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# PowerPoint parsing (fallback)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Excel parsing (fallback)
try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

# Word parsing (fallback)
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


class DocumentParser:
    """Parse various document formats to extract text"""

    def __init__(self, config=None, use_azure_doc_ai=True):
        """
        Initialize document parser

        Args:
            config: Configuration object (optional, will use env vars if not provided)
            use_azure_doc_ai: Whether to use Azure GPT-4o (default: True)
        """
        self.config = config
        self.use_azure_doc_ai = use_azure_doc_ai and HAS_AZURE_DOC_AI
        self.azure_doc_parser = None
        self.docling_parser = None

        # Initialize Azure GPT-4o parser (primary)
        if self.use_azure_doc_ai:
            try:
                self.azure_doc_parser = AzureDocumentParser(config)
                print("✓ Using Azure GPT-4o for document parsing (primary)")
            except Exception as e:
                print(f"⚠ Failed to initialize Azure GPT-4o parser: {e}")
                self.use_azure_doc_ai = False

        # Initialize Docling parser (fallback)
        if HAS_DOCLING:
            try:
                self.docling_parser = DoclingParser()
                print("✓ Using Docling for document parsing (fallback)")
            except Exception as e:
                print(f"⚠ Failed to initialize Docling parser: {e}")

        # Set up supported formats
        self.supported_formats = []
        if self.use_azure_doc_ai and self.azure_doc_parser:
            self.supported_formats = list(self.azure_doc_parser.supported_formats)
        if self.docling_parser:
            for fmt in self.docling_parser.supported_formats:
                if fmt not in self.supported_formats:
                    self.supported_formats.append(fmt)
        # Always add traditional parser formats
        traditional_formats = []
        if HAS_PDF:
            traditional_formats.append('.pdf')
        if HAS_PPTX:
            traditional_formats.append('.pptx')
        if HAS_XLSX:
            traditional_formats.extend(['.xlsx', '.xls', '.xlsm', '.xlsb'])
        if HAS_DOCX:
            traditional_formats.append('.docx')
        for fmt in traditional_formats:
            if fmt not in self.supported_formats:
                self.supported_formats.append(fmt)

    def can_parse(self, file_path: str) -> bool:
        """Check if file format is supported"""
        ext = Path(file_path).suffix.lower()
        return ext in self.supported_formats

    def parse(self, file_path: str) -> Optional[Dict]:
        """
        Parse a document and return extracted text
        Uses Azure Mistral Document AI if available, falls back to traditional parsers

        Returns:
            Dict with 'content' and 'metadata' or None if parsing failed
        """
        if not os.path.exists(file_path):
            return None

        ext = Path(file_path).suffix.lower()

        # Spreadsheet formats: always use traditional parsers (best for tabular data)
        spreadsheet_exts = ('.xlsx', '.xls', '.xlsm', '.xlsb', '.csv', '.tsv', '.ods')

        # Layer 1: Try Azure GPT-4o first for non-spreadsheet files
        if self.use_azure_doc_ai and self.azure_doc_parser and ext not in spreadsheet_exts:
            try:
                result = self.azure_doc_parser.parse(file_path)
                if result and result.get('content') and len(result['content'].strip()) > 10:
                    return result
                print(f"  ⚠ GPT-4o returned insufficient content for {Path(file_path).name}")
            except Exception as e:
                print(f"  ⚠ GPT-4o parser failed: {e}")

        # Layer 2: Try Docling (local, free) for non-spreadsheet files
        if self.docling_parser and ext not in spreadsheet_exts:
            try:
                result = self.docling_parser.parse(file_path)
                if result and result.get('content') and len(result['content'].strip()) > 10:
                    print(f"  ✓ Docling parsed {Path(file_path).name} successfully")
                    return result
                print(f"  ⚠ Docling returned insufficient content for {Path(file_path).name}")
            except Exception as e:
                print(f"  ⚠ Docling parser failed: {e}")

        if ext not in spreadsheet_exts:
            print(f"  Falling back to traditional parser for {Path(file_path).name}")

        # Layer 3: Traditional parsers
        try:
            result = None
            if ext == '.pdf' and HAS_PDF:
                result = self._parse_pdf(file_path)
            elif ext == '.pptx' and HAS_PPTX:
                result = self._parse_pptx(file_path)
            elif ext in ('.xlsx', '.xls', '.xlsm', '.xlsb') and HAS_XLSX:
                result = self._parse_xlsx(file_path)
            elif ext == '.docx' and HAS_DOCX:
                result = self._parse_docx(file_path)

            # Sanitize content - remove NUL characters that break database storage
            if result and result.get('content'):
                result['content'] = result['content'].replace('\x00', '')

            return result
        except Exception as e:
            print(f"  ⚠ Error parsing {Path(file_path).name}: {e}")
            return None

        return None

    def parse_file_bytes(self, file_bytes: bytes, filename: str) -> Optional[str]:
        """
        Parse file bytes using Mistral Document AI (with local fallback).
        Saves bytes to temp file, calls self.parse(), returns content string.

        Args:
            file_bytes: Raw file content as bytes
            filename: Original filename (used to determine file type)

        Returns:
            Extracted text string, or None if parsing failed
        """
        ext = Path(filename).suffix.lower()
        if not ext:
            return None

        MAX_PARSE_BYTES = 50 * 1024 * 1024  # 50 MB general limit
        if len(file_bytes) > MAX_PARSE_BYTES:
            size_mb = len(file_bytes) / (1024 * 1024)
            return f"[File too large to parse: {filename} ({size_mb:.1f} MB). Metadata stored but content not indexed.]"

        # CSV/TSV files - parse as tabular data
        if ext in ('.csv', '.tsv'):
            try:
                return self._parse_csv_bytes(file_bytes, ext)
            except Exception:
                return file_bytes.decode('utf-8', errors='ignore')

        # HTML files - strip tags
        if ext in ('.html', '.htm'):
            return self._parse_html_bytes(file_bytes)

        # JSON files - extract structure
        if ext == '.json':
            return self._parse_json_bytes(file_bytes)

        # XML files - extract text content
        if ext == '.xml':
            return self._parse_xml_bytes(file_bytes)

        # Plain text / markdown / code / config files — read as text
        TEXT_EXTENSIONS = {
            '.txt', '.md', '.rst', '.tex', '.bib',
            # Code
            '.py', '.js', '.ts', '.jsx', '.tsx', '.mjs', '.cjs',
            '.java', '.go', '.rb', '.php', '.cs', '.cpp', '.c', '.h', '.hpp',
            '.rs', '.kt', '.swift', '.scala', '.r', '.rmd',
            '.sh', '.bash', '.zsh', '.ps1', '.bat',
            '.sql', '.graphql', '.proto',
            '.css', '.scss', '.sass', '.less',
            '.vue', '.svelte', '.ipynb',
            '.d.ts', '.js.map', '.css.map',
            # Config
            '.yaml', '.yml', '.toml', '.ini', '.conf', '.cfg', '.env',
            '.lock', '.mod', '.sum',
            '.log', '.ndjson', '.jsonl', '.geojson',
        }
        if ext in TEXT_EXTENSIONS:
            try:
                return file_bytes.decode('utf-8', errors='ignore')
            except Exception:
                return None

        # RTF files
        if ext == '.rtf':
            return self._parse_rtf_bytes(file_bytes)

        # Image files - describe via GPT-4o vision
        if ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.tif', '.tiff'):
            if len(file_bytes) > 20 * 1024 * 1024:
                return f"[Image: {filename} — too large for vision analysis ({len(file_bytes) // (1024*1024)} MB)]"
            return self._parse_image_bytes(file_bytes, filename)

        # Audio/video files - transcribe via Whisper
        audio_exts = ('.mp3', '.wav', '.m4a', '.webm', '.mp4', '.mov', '.ogg', '.flac')
        if ext in audio_exts:
            return self._parse_audio_bytes(file_bytes, filename)

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            result = self.parse(tmp_path)
            if result and result.get('content'):
                return result['content']
            # Return a minimal description for empty files rather than None
            if ext in ('.xlsx', '.xls', '.xlsm', '.xlsb', '.csv', '.tsv'):
                return f'[Empty spreadsheet: {filename}]'
            return None
        except Exception as e:
            print(f"[DocumentParser] parse_file_bytes error for {filename}: {e}")
            return None
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def _parse_pdf(self, file_path: str) -> Dict:
        """Extract text from PDF"""
        text_parts = []

        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            MAX_PDF_PAGES = 500
            total_pages = len(pdf_reader.pages)
            num_pages = min(total_pages, MAX_PDF_PAGES)
            if total_pages > MAX_PDF_PAGES:
                text_parts.append(f"[NOTE: PDF has {total_pages} pages. Only first {MAX_PDF_PAGES} indexed.]")

            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text and text.strip():
                    # Remove NUL characters that cause database issues
                    clean_text = text.replace('\x00', '').strip()
                    if clean_text:
                        text_parts.append(clean_text)

        content = '\n\n'.join(text_parts)
        # Final cleanup of any remaining NUL characters
        content = content.replace('\x00', '')

        return {
            'content': content,
            'metadata': {
                'pages': total_pages,
                'pages_indexed': num_pages,
                'file_type': 'pdf'
            }
        }

    def _parse_pptx(self, file_path: str) -> Dict:
        """Extract text from PowerPoint"""
        text_parts = []

        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'slides': len(prs.slides),
                'file_type': 'pptx'
            }
        }

    def _format_tabular_data(self, rows: List[List[str]], sheet_name: str, max_rows: int = 10000) -> str:
        """Format tabular data into AI-readable labeled-row format for RAG retrieval.

        Args:
            rows: List of rows, each row is a list of cell values as strings.
                  Row 0 is treated as headers if it looks like a header row.
            sheet_name: Name of the sheet/table
            max_rows: Maximum data rows to include

        Returns:
            Formatted string with column context per row
        """
        if not rows:
            return ''

        # Trim trailing empty columns
        max_cols = 0
        for row in rows:
            last_non_empty = -1
            for i, cell in enumerate(row):
                if cell.strip():
                    last_non_empty = i
            if last_non_empty + 1 > max_cols:
                max_cols = last_non_empty + 1
        if max_cols == 0:
            return ''
        rows = [row[:max_cols] for row in rows]
        # Pad short rows
        rows = [row + [''] * (max_cols - len(row)) for row in rows]

        # Detect if row 0 is a header row (non-empty, non-numeric strings)
        first_row = rows[0]
        is_header = all(
            cell.strip() and not cell.strip().replace('.', '').replace(',', '').replace('-', '').isdigit()
            for cell in first_row if cell.strip()
        ) and any(cell.strip() for cell in first_row)

        if is_header:
            headers = [cell.strip() or f'Column {chr(65 + i)}' for i, cell in enumerate(first_row)]
            data_rows = rows[1:]
        else:
            headers = [f'Column {chr(65 + i)}' if i < 26 else f'Column {i + 1}' for i in range(max_cols)]
            data_rows = rows

        # Limit data rows
        truncated = len(data_rows) > max_rows
        data_rows = data_rows[:max_rows]

        # Filter out fully empty data rows
        data_rows = [r for r in data_rows if any(cell.strip() for cell in r)]

        num_cols = len(headers)
        num_data_rows = len(data_rows)

        parts = []
        parts.append(f'[Sheet: {sheet_name}]')
        parts.append(f'Columns: {" | ".join(headers)}')
        parts.append(f'Rows: {num_data_rows}')
        parts.append('')

        # For wide sheets (>20 cols), use compact markdown table
        if num_cols > 20:
            parts.append('| ' + ' | '.join(headers) + ' |')
            parts.append('| ' + ' | '.join(['---'] * num_cols) + ' |')
            for row in data_rows:
                parts.append('| ' + ' | '.join(cell.strip() for cell in row) + ' |')
        else:
            # Labeled-row format (best for RAG)
            for i, row in enumerate(data_rows, 1):
                row_lines = [f'Row {i}:']
                for j, cell in enumerate(row):
                    val = cell.strip()
                    if val:
                        row_lines.append(f'- {headers[j]}: {val}')
                if len(row_lines) > 1:  # Has at least one non-empty value
                    parts.append('\n'.join(row_lines))
                    parts.append('')

        if truncated:
            parts.append(f'\n[WARNING] Sheet truncated at {max_rows:,} rows.')

        return '\n'.join(parts)

    def _parse_csv_bytes(self, file_bytes: bytes, ext: str = '.csv') -> Optional[str]:
        """Parse CSV/TSV bytes into AI-readable tabular format"""
        MAX_ROWS = 10000

        # Decode with BOM handling
        try:
            text = file_bytes.decode('utf-8-sig')
        except UnicodeDecodeError:
            text = file_bytes.decode('latin-1', errors='ignore')

        delimiter = '\t' if ext == '.tsv' else ','

        try:
            # Sniff the dialect for better delimiter detection
            sample = text[:8192]
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=',\t;|')
                reader = csv.reader(io.StringIO(text), dialect)
            except csv.Error:
                reader = csv.reader(io.StringIO(text), delimiter=delimiter)

            rows = []
            for i, row in enumerate(reader):
                if i >= MAX_ROWS + 1:  # +1 for header
                    break
                rows.append([cell.strip() for cell in row])

            if not rows:
                return file_bytes.decode('utf-8', errors='ignore')

            return self._format_tabular_data(rows, 'Data')
        except Exception:
            # Fallback: return raw text
            return text

    def _parse_xlsx(self, file_path: str) -> Dict:
        """Extract text from Excel files (.xlsx, .xls, .xlsm, .xlsb) with 10K row limit per sheet"""
        MAX_ROWS_PER_SHEET = 10000
        ext = Path(file_path).suffix.lower()

        # For .xls and .xlsb, use pandas (openpyxl can't read these)
        if ext in ('.xls', '.xlsb'):
            return self._parse_excel_pandas(file_path, MAX_ROWS_PER_SHEET)

        # For .xlsx and .xlsm, try openpyxl first, fall back to pandas
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        except Exception as e:
            print(f"[DocumentParser] openpyxl failed for {Path(file_path).name}: {e}, falling back to pandas")
            return self._parse_excel_pandas(file_path, MAX_ROWS_PER_SHEET)

        text_parts = []
        total_rows = 0
        truncated_sheets = []
        sheet_names = list(wb.sheetnames)  # Save before close

        try:
            for sheet_name in sheet_names:
                sheet = wb[sheet_name]
                rows_collected = []

                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= MAX_ROWS_PER_SHEET:
                        truncated_sheets.append(sheet_name)
                        break

                    # Preserve all cells including empty ones to maintain column alignment
                    row_values = [str(cell).strip() if cell is not None else '' for cell in row]
                    # Skip fully empty rows
                    if any(v for v in row_values):
                        rows_collected.append(row_values)
                        row_count += 1

                if rows_collected:
                    text_parts.append(self._format_tabular_data(rows_collected, sheet_name, MAX_ROWS_PER_SHEET))
                    total_rows += row_count
        except Exception as e:
            print(f"[DocumentParser] openpyxl iteration failed: {e}, falling back to pandas")
            try:
                wb.close()
            except Exception:
                pass
            return self._parse_excel_pandas(file_path, MAX_ROWS_PER_SHEET)

        wb.close()

        if truncated_sheets:
            warning = f"\n\n[WARNING] The following sheets exceeded {MAX_ROWS_PER_SHEET:,} rows and were truncated: {', '.join(truncated_sheets)}"
            text_parts.append(warning)

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'sheets': len(sheet_names),
                'total_rows': total_rows,
                'truncated_sheets': truncated_sheets,
                'max_rows_per_sheet': MAX_ROWS_PER_SHEET,
                'file_type': ext.lstrip('.')
            }
        }

    def _parse_excel_pandas(self, file_path: str, max_rows: int) -> Dict:
        """Fallback Excel parser using pandas for .xls and .xlsb formats"""
        import pandas as pd

        text_parts = []
        total_rows = 0
        truncated_sheets = []

        try:
            xls = pd.ExcelFile(file_path)
            sheet_names = xls.sheet_names

            for sheet_name in sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=max_rows)
                if len(df) >= max_rows:
                    truncated_sheets.append(sheet_name)

                # Collect rows preserving column alignment
                rows_collected = []
                for _, row in df.iterrows():
                    row_values = [str(v).strip() if pd.notna(v) else '' for v in row]
                    if any(v for v in row_values):
                        rows_collected.append(row_values)
                        total_rows += 1

                if rows_collected:
                    text_parts.append(self._format_tabular_data(rows_collected, sheet_name, max_rows))

            xls.close()
        except Exception as e:
            print(f"[DocumentParser] pandas Excel parse error: {e}")
            return {'content': '', 'metadata': {'error': str(e), 'file_type': Path(file_path).suffix.lstrip('.')}}

        if truncated_sheets:
            warning = f"\n\n[WARNING] The following sheets exceeded {max_rows:,} rows and were truncated: {', '.join(truncated_sheets)}"
            text_parts.append(warning)

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'sheets': len(sheet_names),
                'total_rows': total_rows,
                'truncated_sheets': truncated_sheets,
                'max_rows_per_sheet': max_rows,
                'file_type': Path(file_path).suffix.lstrip('.')
            }
        }

    def _parse_docx(self, file_path: str) -> Dict:
        """Extract text from Word document"""
        doc = Document(file_path)

        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        content = '\n\n'.join(text_parts)

        return {
            'content': content,
            'metadata': {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'file_type': 'docx'
            }
        }

    def _parse_html_bytes(self, file_bytes):
        """Extract clean text from HTML"""
        try:
            from bs4 import BeautifulSoup
            html = file_bytes.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'html.parser')
            for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
                tag.decompose()
            text = soup.get_text(separator='\n', strip=True)
            import re
            text = re.sub(r'\n{3,}', '\n\n', text)
            return text
        except Exception:
            return file_bytes.decode('utf-8', errors='ignore')

    def _parse_json_bytes(self, file_bytes):
        """Extract readable content from JSON"""
        import json as json_mod
        try:
            data = json_mod.loads(file_bytes.decode('utf-8', errors='ignore'))
            return self._json_to_text(data, max_depth=5)
        except (json_mod.JSONDecodeError, Exception):
            return file_bytes.decode('utf-8', errors='ignore')

    def _json_to_text(self, data, max_depth=5, depth=0, prefix=''):
        """Convert JSON structure to readable text"""
        if depth > max_depth:
            return '[nested data truncated]'
        parts = []
        if isinstance(data, dict):
            for key, value in list(data.items())[:100]:
                if isinstance(value, (dict, list)):
                    parts.append(f"{prefix}{key}:")
                    parts.append(self._json_to_text(value, max_depth, depth + 1, prefix + '  '))
                else:
                    parts.append(f"{prefix}{key}: {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data[:200]):
                parts.append(f"{prefix}[{i}] {self._json_to_text(item, max_depth, depth + 1, prefix + '  ')}")
        else:
            return str(data)
        return '\n'.join(parts)

    def _parse_xml_bytes(self, file_bytes):
        """Extract text content from XML"""
        try:
            import xml.etree.ElementTree as ET
            root = ET.fromstring(file_bytes.decode('utf-8', errors='ignore'))
            texts = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                    texts.append(f"{tag}: {elem.text.strip()}")
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            return '\n'.join(texts[:5000])
        except Exception:
            return file_bytes.decode('utf-8', errors='ignore')

    def _parse_rtf_bytes(self, file_bytes):
        """Extract text from RTF files"""
        try:
            from striprtf.striprtf import rtf_to_text
            rtf_content = file_bytes.decode('utf-8', errors='ignore')
            return rtf_to_text(rtf_content)
        except ImportError:
            import re
            text = file_bytes.decode('utf-8', errors='ignore')
            text = re.sub(r'\\[a-z]+\d*\s?', '', text)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            return text if len(text) > 20 else None
        except Exception:
            return None

    def _parse_image_bytes(self, file_bytes, filename):
        """Describe image content using GPT-4o vision"""
        import base64
        import io as _io
        try:
            from openai import AzureOpenAI
            client = AzureOpenAI(
                azure_endpoint=os.environ.get('AZURE_OPENAI_ENDPOINT'),
                api_key=os.environ.get('AZURE_OPENAI_API_KEY'),
                api_version="2024-12-01-preview"
            )
            ext = Path(filename).suffix.lower().lstrip('.')

            # Azure OpenAI vision only supports PNG, JPEG, GIF, WebP
            # Convert BMP/TIFF/other formats to PNG first
            if ext not in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(_io.BytesIO(file_bytes))
                    buf = _io.BytesIO()
                    img.save(buf, format="PNG")
                    file_bytes = buf.getvalue()
                    ext = 'png'
                except Exception as conv_err:
                    print(f"[DocumentParser] Image conversion failed for {filename}: {conv_err}")
                    return f"[Image file: {filename} — format conversion failed]"

            b64 = base64.b64encode(file_bytes).decode('utf-8')
            mime_map = {'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'gif': 'image/gif', 'webp': 'image/webp'}
            mime = mime_map.get(ext, 'image/png')
            response = client.chat.completions.create(
                model=os.environ.get('AZURE_CHAT_DEPLOYMENT', 'gpt-5-chat'),
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in detail. Include any text, labels, data, charts, diagrams, or notable content. Be thorough and specific."},
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}
                    ]
                }],
                max_tokens=1000
            )
            description = response.choices[0].message.content
            return f"[Image: {filename}]\n{description}"
        except Exception as e:
            print(f"[DocumentParser] Image parsing error for {filename}: {e}")
            return f"[Image file: {filename} — could not extract content]"

    def _parse_audio_bytes(self, file_bytes, filename):
        """Transcribe audio using Azure Whisper"""
        ext = Path(filename).suffix.lower()
        try:
            from openai import AzureOpenAI
            client = AzureOpenAI(
                azure_endpoint=os.environ.get('AZURE_OPENAI_ENDPOINT'),
                api_key=os.environ.get('AZURE_OPENAI_API_KEY'),
                api_version="2024-12-01-preview"
            )
            if len(file_bytes) > 25 * 1024 * 1024:
                return f"[Audio: {filename} — too large for transcription ({len(file_bytes) // (1024*1024)} MB, limit 25 MB)]"

            tmp_path = None
            try:
                with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                with open(tmp_path, 'rb') as audio_file:
                    transcript = client.audio.transcriptions.create(
                        model=os.environ.get('AZURE_WHISPER_DEPLOYMENT', 'whisper'),
                        file=audio_file,
                        response_format="text"
                    )
                return f"[Audio transcription: {filename}]\n{transcript}"
            finally:
                if tmp_path and os.path.exists(tmp_path):
                    os.unlink(tmp_path)
        except Exception as e:
            print(f"[DocumentParser] Audio transcription error for {filename}: {e}")
            return f"[Audio file: {filename} — transcription failed]"

    def parse_pdf_bytes(self, content: bytes) -> str:
        """
        Extract text from PDF bytes (for file uploads)

        Args:
            content: PDF file content as bytes

        Returns:
            Extracted text string
        """
        import io

        # Try Azure Document AI first if available
        if HAS_AZURE_DOC_AI and self.azure_doc_parser:
            try:
                result = self.azure_doc_parser.parse_bytes(content, 'application/pdf')
                if result.get('success') and result.get('content'):
                    return result['content']
            except Exception as e:
                print(f"[DocumentParser] Azure Document AI failed, falling back to PyPDF2: {e}")

        # Fall back to PyPDF2
        if not HAS_PDF:
            raise Exception("PDF parsing not available. PyPDF2 is not installed.")

        text_parts = []
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            num_pages = len(pdf_reader.pages)

            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text and text.strip():
                    # Remove NUL characters that cause database issues
                    clean_text = text.replace('\x00', '').strip()
                    if clean_text:
                        text_parts.append(clean_text)

            result = '\n\n'.join(text_parts)
            # Final cleanup of any remaining NUL characters
            return result.replace('\x00', '')
        except Exception as e:
            raise Exception(f"Failed to parse PDF: {str(e)}")

    def parse_word_bytes(self, content: bytes) -> str:
        """
        Extract text from Word document bytes (for file uploads)

        Args:
            content: DOCX file content as bytes

        Returns:
            Extracted text string
        """
        import io

        # Try Azure Document AI first if available
        if HAS_AZURE_DOC_AI and self.azure_doc_parser:
            try:
                result = self.azure_doc_parser.parse_bytes(content, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                if result.get('success') and result.get('content'):
                    return result['content']
            except Exception as e:
                print(f"[DocumentParser] Azure Document AI failed, falling back to python-docx: {e}")

        # Fall back to python-docx
        if not HAS_DOCX:
            raise Exception("Word document parsing not available. python-docx is not installed.")

        text_parts = []
        try:
            doc = Document(io.BytesIO(content))

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text_parts.append(row_text)

            return '\n\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Failed to parse Word document: {str(e)}")


if __name__ == "__main__":
    # Test the parser
    parser = DocumentParser()
    print(f"Supported formats: {parser.supported_formats}")

    # Test with a sample file
    test_file = "/Users/rishitjain/Downloads/Takeout/Google Chat/Groups/Space AAAAn7sv4eE/File-Timeline - BEAT Healthcare Consulting.pptx"
    if os.path.exists(test_file):
        result = parser.parse(test_file)
        if result:
            print(f"\nExtracted {len(result['content'])} characters")
            print(f"Preview: {result['content'][:200]}...")
