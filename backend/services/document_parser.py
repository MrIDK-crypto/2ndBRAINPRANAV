"""
Universal Document Parser Service
Priority: Local parsing (instant) → GPT-4o vision (fast) → LlamaParse (fallback)

- DOCX/PPTX/XLSX: python-docx/python-pptx/openpyxl (instant, no API)
- PDF: PyPDF2 (instant for text PDFs) → GPT-4o vision (for scanned)
- Images: GPT-4o vision
- Fallback: LlamaParse API
"""

import os
import io
import base64
import time
import httpx
import asyncio
from typing import Optional, Dict, Any
from pathlib import Path
from dotenv import load_dotenv

load_dotenv(override=True)

# Local parsers (instant, no API calls)
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from striprtf.striprtf import rtf_to_text
    HAS_RTF = True
except ImportError:
    HAS_RTF = False

# Azure OpenAI for GPT-4o vision (scanned PDFs, images)
try:
    from openai import AzureOpenAI
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False

# LlamaParse configuration (last resort fallback)
def _get_llama_key():
    return os.getenv("LLAMA_CLOUD_API_KEY", "")

LLAMAPARSE_API_URL = "https://api.cloud.llamaindex.ai/api/parsing"

# Azure Document Intelligence configuration
def _get_azure_config():
    return (
        os.getenv("AZURE_DI_ENDPOINT", ""),
        os.getenv("AZURE_DI_API_KEY", "")
    )

# Minimum chars to consider local extraction successful
MIN_CONTENT_CHARS = 50


class DocumentParser:
    """
    Universal document parser with speed-first routing:
    1. Local parsers (instant) for DOCX, PPTX, XLSX, PDF, RTF
    2. GPT-4o vision for scanned PDFs and images
    3. LlamaParse API as last resort fallback
    """

    # Extensions that can be parsed locally (instant)
    LOCAL_PARSE_EXTENSIONS = {
        ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".rtf",
    }

    # Image extensions → GPT-4o vision
    IMAGE_EXTENSIONS = {
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif",
    }

    # Plain text files (read directly) — includes code, config, data formats
    PLAIN_TEXT_EXTENSIONS = {
        ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".r", ".rmd",
        # Code
        ".py", ".js", ".ts", ".jsx", ".tsx", ".mjs", ".cjs",
        ".java", ".go", ".rb", ".php", ".cs", ".cpp", ".c", ".h", ".hpp",
        ".rs", ".kt", ".swift", ".scala",
        ".sh", ".bash", ".zsh", ".ps1", ".bat",
        ".sql", ".graphql", ".proto",
        ".css", ".scss", ".sass", ".less",
        ".vue", ".svelte", ".ipynb", ".tex", ".bib", ".rst",
        ".d.ts", ".js.map", ".css.map",
        # Config
        ".yaml", ".yml", ".toml", ".ini", ".conf", ".cfg", ".env",
        ".lock", ".mod", ".sum",
        ".log", ".tsv", ".ndjson", ".jsonl", ".geojson",
    }

    # LlamaParse fallback extensions (ODF formats without local parsers)
    LLAMAPARSE_EXTENSIONS = {".odt", ".ods", ".odp"}

    # All supported extensions
    SUPPORTED_EXTENSIONS = LOCAL_PARSE_EXTENSIONS | IMAGE_EXTENSIONS | PLAIN_TEXT_EXTENSIONS | LLAMAPARSE_EXTENSIONS

    def __init__(self, llama_api_key: Optional[str] = None, azure_endpoint: Optional[str] = None, azure_api_key: Optional[str] = None):
        # Get fresh values from environment
        azure_ep, azure_key = _get_azure_config()

        self.llama_api_key = llama_api_key or _get_llama_key()
        self.azure_endpoint = (azure_endpoint or azure_ep).rstrip('/')
        self.azure_api_key = azure_api_key or azure_key

        # Initialize GPT-4o client for vision tasks
        self.openai_client = None
        self.chat_model = None
        if HAS_OPENAI:
            try:
                oai_key = os.getenv("AZURE_OPENAI_API_KEY", "")
                oai_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", "")
                oai_version = os.getenv("AZURE_API_VERSION", "2024-12-01-preview")
                if oai_key and oai_endpoint:
                    self.openai_client = AzureOpenAI(
                        azure_endpoint=oai_endpoint,
                        api_key=oai_key,
                        api_version=oai_version,
                    )
                    self.chat_model = os.getenv("AZURE_CHAT_DEPLOYMENT", "gpt-5-chat")
                    print(f"[DocumentParser] GPT-4o vision ready (model: {self.chat_model})")
            except Exception as e:
                print(f"[DocumentParser] GPT-4o init failed: {e}")

        caps = []
        if HAS_DOCX: caps.append("DOCX")
        if HAS_PPTX: caps.append("PPTX")
        if HAS_XLSX: caps.append("XLSX")
        if HAS_PDF: caps.append("PDF")
        if HAS_RTF: caps.append("RTF")
        if self.openai_client: caps.append("GPT-4o-vision")
        print(f"[DocumentParser] Local parsers: {', '.join(caps) or 'none'}")

    def is_supported(self, file_extension: str) -> bool:
        ext = file_extension.lower() if file_extension.startswith('.') else f".{file_extension.lower()}"
        return ext in self.SUPPORTED_EXTENSIONS

    def is_plain_text(self, file_extension: str) -> bool:
        ext = file_extension.lower() if file_extension.startswith('.') else f".{file_extension.lower()}"
        return ext in self.PLAIN_TEXT_EXTENSIONS

    # ─── Main entry point ───────────────────────────────────────────

    async def parse_bytes(
        self,
        file_bytes: bytes,
        file_name: str,
        file_extension: str
    ) -> str:
        ext = file_extension.lower() if file_extension.startswith('.') else f".{file_extension.lower()}"

        # Plain text — decode directly
        if self.is_plain_text(ext):
            try:
                return file_bytes.decode('utf-8', errors='ignore')
            except Exception as e:
                print(f"[DocumentParser] Error decoding text file: {e}")
                return ""

        # Images — GPT-4o vision
        if ext in self.IMAGE_EXTENSIONS:
            return await self._parse_image_with_gpt4o(file_bytes, file_name, ext)

        # Office docs & PDFs — try local first, then GPT-4o, then LlamaParse
        if ext in self.LOCAL_PARSE_EXTENSIONS:
            # Step 1: Local parsing (instant)
            content = self._parse_locally(file_bytes, file_name, ext)
            if content and len(content.strip()) >= MIN_CONTENT_CHARS:
                print(f"[DocumentParser] Local parse OK: {file_name} → {len(content)} chars")
                return content

            # Step 2: For PDFs with no/minimal text (scanned), try GPT-4o vision
            if ext == ".pdf" and self.openai_client:
                print(f"[DocumentParser] Local PDF extraction minimal, trying GPT-4o vision for {file_name}")
                gpt_content = await self._parse_pdf_with_gpt4o(file_bytes, file_name)
                if gpt_content and len(gpt_content.strip()) >= MIN_CONTENT_CHARS:
                    return gpt_content

            # Step 3: LlamaParse fallback
            if self.llama_api_key:
                print(f"[DocumentParser] Falling back to LlamaParse for {file_name}")
                return await self._parse_with_llamaparse(file_bytes, file_name, ext)

            # Return whatever local parsing got (even if minimal)
            return content or ""

        # ODF formats — LlamaParse only
        if ext in self.LLAMAPARSE_EXTENSIONS:
            if self.llama_api_key:
                return await self._parse_with_llamaparse(file_bytes, file_name, ext)
            print(f"[DocumentParser] No parser available for {ext}")
            return ""

        print(f"[DocumentParser] Unsupported file type: {ext}")
        return ""

    async def parse_file(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            print(f"[DocumentParser] File not found: {file_path}")
            return ""
        ext = path.suffix.lower()
        with open(file_path, 'rb') as f:
            file_bytes = f.read()
        return await self.parse_bytes(file_bytes, path.name, ext)

    # ─── Local parsers (instant, no API calls) ──────────────────────

    def _parse_locally(self, file_bytes: bytes, file_name: str, ext: str) -> str:
        """Parse document using local libraries. Returns extracted text or empty string."""
        try:
            if ext == ".docx" and HAS_DOCX:
                return self._parse_docx(file_bytes)
            elif ext == ".doc":
                # .doc (legacy) - no good Python parser, skip to fallback
                return ""
            elif ext == ".pptx" and HAS_PPTX:
                return self._parse_pptx(file_bytes)
            elif ext == ".ppt":
                return ""
            elif ext == ".xlsx" and HAS_XLSX:
                return self._parse_xlsx(file_bytes)
            elif ext == ".xls":
                return ""
            elif ext == ".pdf" and HAS_PDF:
                return self._parse_pdf(file_bytes)
            elif ext == ".rtf" and HAS_RTF:
                return self._parse_rtf(file_bytes)
        except Exception as e:
            print(f"[DocumentParser] Local parse error for {file_name}: {e}")
        return ""

    def _parse_docx(self, file_bytes: bytes) -> str:
        doc = DocxDocument(io.BytesIO(file_bytes))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return '\n\n'.join(parts)

    def _parse_pptx(self, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + '\n'.join(slide_texts))
        return '\n\n'.join(parts)

    def _parse_xlsx(self, file_bytes: bytes) -> str:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        parts = []
        MAX_ROWS = 10000
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts = [f"[Sheet: {sheet_name}]"]
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > MAX_ROWS:
                    sheet_parts.append(f"... (truncated at {MAX_ROWS} rows)")
                    break
                cells = [str(c) if c is not None else '' for c in row]
                line = ' | '.join(cells)
                if line.strip() and line.strip(' |'):
                    sheet_parts.append(line)
            parts.append('\n'.join(sheet_parts))
        wb.close()
        return '\n\n'.join(parts)

    def _parse_pdf(self, file_bytes: bytes) -> str:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.replace('\x00', '').strip())
        return '\n\n'.join(parts)

    def _parse_rtf(self, file_bytes: bytes) -> str:
        text = file_bytes.decode('utf-8', errors='ignore')
        return rtf_to_text(text)

    # ─── GPT-4o vision (for scanned PDFs and images) ────────────────

    async def _parse_image_with_gpt4o(self, file_bytes: bytes, file_name: str, ext: str) -> str:
        """Parse image using GPT-4o vision."""
        if not self.openai_client:
            # Fall back to Azure Document Intelligence for images
            if self.azure_endpoint and self.azure_api_key:
                return await self._parse_with_azure_di(file_bytes, file_name, ext)
            print(f"[DocumentParser] No vision parser available for {file_name}")
            return ""

        try:
            print(f"[DocumentParser] Parsing {file_name} with GPT-4o vision")

            # Azure OpenAI vision only supports PNG, JPEG, GIF, WebP
            # Convert BMP/TIFF/other formats to PNG first
            NATIVE_FORMATS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
            if ext not in NATIVE_FORMATS:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(io.BytesIO(file_bytes))
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    file_bytes = buf.getvalue()
                    ext = ".png"
                    print(f"[DocumentParser] Converted {file_name} to PNG for GPT-4o vision")
                except Exception as conv_err:
                    print(f"[DocumentParser] Image conversion failed for {file_name}: {conv_err}")
                    return ""

            b64 = base64.b64encode(file_bytes).decode('utf-8')
            mime_map = {
                ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp",
            }
            mime = mime_map.get(ext, "image/png")

            response = self.openai_client.chat.completions.create(
                model=self.chat_model,
                messages=[
                    {"role": "system", "content": "Extract ALL text and content from this image. Return clean text preserving structure. Include tables in markdown format. Do NOT add commentary."},
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract all content from this image:"},
                            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"}}
                        ]
                    }
                ],
                temperature=0.1,
                max_tokens=16000,
            )
            content = response.choices[0].message.content
            if content and len(content.strip()) >= 10:
                print(f"[DocumentParser] GPT-4o vision extracted {len(content)} chars from {file_name}")
                return content
        except Exception as e:
            print(f"[DocumentParser] GPT-4o vision error for {file_name}: {e}")

        # Fall back to Azure DI
        if self.azure_endpoint and self.azure_api_key:
            return await self._parse_with_azure_di(file_bytes, file_name, ext)
        return ""

    async def _parse_pdf_with_gpt4o(self, file_bytes: bytes, file_name: str) -> str:
        """Parse scanned PDF by converting pages to images and using GPT-4o vision."""
        if not self.openai_client:
            return ""

        try:
            # Try PyMuPDF (fitz) for page-to-image conversion
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            all_text = []

            # Process up to 20 pages to avoid excessive API calls
            max_pages = min(len(doc), 20)
            print(f"[DocumentParser] Converting {max_pages} PDF pages to images for GPT-4o ({file_name})")

            for page_num in range(max_pages):
                page = doc[page_num]
                # Render page to image at 150 DPI
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                b64 = base64.b64encode(img_bytes).decode('utf-8')

                response = self.openai_client.chat.completions.create(
                    model=self.chat_model,
                    messages=[
                        {"role": "system", "content": "Extract ALL text from this PDF page image. Preserve structure, tables, and formatting. Return only the extracted content."},
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": f"Extract all text from page {page_num + 1}:"},
                                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"}}
                            ]
                        }
                    ],
                    temperature=0.1,
                    max_tokens=16000,
                )
                page_text = response.choices[0].message.content
                if page_text and page_text.strip():
                    all_text.append(page_text.strip())

            doc.close()
            if all_text:
                content = '\n\n'.join(all_text)
                print(f"[DocumentParser] GPT-4o vision extracted {len(content)} chars from {max_pages} PDF pages")
                return content

        except ImportError:
            print(f"[DocumentParser] PyMuPDF not available, cannot use GPT-4o for scanned PDF")
        except Exception as e:
            print(f"[DocumentParser] GPT-4o PDF parsing error for {file_name}: {e}")

        return ""

    # ─── Azure Document Intelligence (images) ───────────────────────

    async def _parse_with_azure_di(
        self,
        file_bytes: bytes,
        file_name: str,
        extension: str
    ) -> str:
        try:
            print(f"[DocumentParser] Parsing {file_name} ({len(file_bytes)} bytes) with Azure Document Intelligence")

            analyze_url = f"{self.azure_endpoint}/documentintelligence/documentModels/prebuilt-layout:analyze?api-version=2024-11-30"

            headers = {
                "Ocp-Apim-Subscription-Key": self.azure_api_key,
                "Content-Type": self._get_mime_type(extension.lstrip('.'))
            }

            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(analyze_url, content=file_bytes, headers=headers)
                if response.status_code != 202:
                    print(f"[DocumentParser] Azure DI submit failed: {response.status_code} - {response.text}")
                    return ""

                operation_location = response.headers.get("Operation-Location")
                if not operation_location:
                    print("[DocumentParser] No operation location in Azure DI response")
                    return ""

                poll_headers = {"Ocp-Apim-Subscription-Key": self.azure_api_key}
                for attempt in range(60):
                    await asyncio.sleep(2)
                    status_response = await client.get(operation_location, headers=poll_headers)
                    if status_response.status_code != 200:
                        continue
                    result = status_response.json()
                    status = result.get("status")
                    if status == "succeeded":
                        analyze_result = result.get("analyzeResult", {})
                        content = analyze_result.get("content", "")
                        if content:
                            print(f"[DocumentParser] Azure DI: {len(content)} chars from {file_name}")
                            return content
                        paragraphs = analyze_result.get("paragraphs", [])
                        if paragraphs:
                            return "\n\n".join([p.get("content", "") for p in paragraphs])
                        return ""
                    elif status == "failed":
                        error = result.get("error", {})
                        print(f"[DocumentParser] Azure DI failed: {error.get('message', 'Unknown')}")
                        return ""
                    elif status in ["notStarted", "running"]:
                        continue

                print(f"[DocumentParser] Azure DI timeout for {file_name}")
                return ""

        except Exception as e:
            print(f"[DocumentParser] Azure DI error: {e}")
            return ""

    # ─── LlamaParse (last resort fallback) ───────────────────────────

    async def _parse_with_llamaparse(
        self,
        file_bytes: bytes,
        file_name: str,
        extension: str
    ) -> str:
        try:
            ext = extension.lstrip('.')
            print(f"[DocumentParser] Parsing {file_name} ({len(file_bytes)} bytes) with LlamaParse (fallback)")

            async with httpx.AsyncClient(timeout=120.0) as client:
                files = {"file": (file_name, file_bytes, self._get_mime_type(ext))}
                data = {"gpt4o_mode": "true"}
                headers = {"Authorization": f"Bearer {self.llama_api_key}"}

                upload_response = await client.post(
                    f"{LLAMAPARSE_API_URL}/upload",
                    files=files, data=data, headers=headers
                )
                if upload_response.status_code != 200:
                    print(f"[DocumentParser] LlamaParse upload failed: {upload_response.status_code}")
                    return ""

                job_id = upload_response.json().get("id")
                if not job_id:
                    return ""

                for attempt in range(60):
                    status_response = await client.get(
                        f"{LLAMAPARSE_API_URL}/job/{job_id}", headers=headers
                    )
                    if status_response.status_code != 200:
                        await asyncio.sleep(5)
                        continue
                    status_data = status_response.json()
                    status = status_data.get("status")
                    if status == "SUCCESS":
                        break
                    elif status == "ERROR":
                        return ""
                    await asyncio.sleep(5)
                else:
                    return ""

                result_response = await client.get(
                    f"{LLAMAPARSE_API_URL}/job/{job_id}/result/text", headers=headers
                )
                if result_response.status_code != 200:
                    result_response = await client.get(
                        f"{LLAMAPARSE_API_URL}/job/{job_id}/result/markdown", headers=headers
                    )
                if result_response.status_code == 200:
                    try:
                        rj = result_response.json()
                        if isinstance(rj, dict):
                            text = rj.get("text") or rj.get("markdown") or result_response.text
                        else:
                            text = result_response.text
                    except Exception:
                        text = result_response.text
                    print(f"[DocumentParser] LlamaParse: {len(text)} chars from {file_name}")
                    return text
                return ""

        except Exception as e:
            print(f"[DocumentParser] LlamaParse error: {e}")
            return ""

    # ─── Synchronous versions (for gevent/gunicorn workers) ────────

    def parse_bytes_sync(self, file_bytes: bytes, file_name: str, file_extension: str) -> str:
        """Fully synchronous parse — no asyncio. Safe in gevent workers."""
        ext = file_extension.lower() if file_extension.startswith('.') else f".{file_extension.lower()}"

        if self.is_plain_text(ext):
            try:
                return file_bytes.decode('utf-8', errors='ignore')
            except Exception:
                return ""

        if ext in self.IMAGE_EXTENSIONS:
            return self._parse_image_with_gpt4o_sync(file_bytes, file_name, ext)

        if ext in self.LOCAL_PARSE_EXTENSIONS:
            content = self._parse_locally(file_bytes, file_name, ext)
            if content and len(content.strip()) >= MIN_CONTENT_CHARS:
                return content
            if ext == ".pdf" and self.openai_client:
                gpt_content = self._parse_pdf_with_gpt4o_sync(file_bytes, file_name)
                if gpt_content and len(gpt_content.strip()) >= MIN_CONTENT_CHARS:
                    return gpt_content
            if self.llama_api_key:
                return self._parse_with_llamaparse_sync(file_bytes, file_name, ext)
            return content or ""

        if ext in self.LLAMAPARSE_EXTENSIONS:
            if self.llama_api_key:
                return self._parse_with_llamaparse_sync(file_bytes, file_name, ext)
            return ""

        return ""

    def _parse_image_with_gpt4o_sync(self, file_bytes: bytes, file_name: str, ext: str) -> str:
        """Synchronous GPT-4o vision image parsing."""
        if not self.openai_client:
            if self.azure_endpoint and self.azure_api_key:
                return self._parse_with_azure_di_sync(file_bytes, file_name, ext)
            return ""
        try:
            NATIVE_FORMATS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
            if ext not in NATIVE_FORMATS:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(io.BytesIO(file_bytes))
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    file_bytes = buf.getvalue()
                    ext = ".png"
                except Exception:
                    return ""

            b64 = base64.b64encode(file_bytes).decode('utf-8')
            mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".gif": "image/gif", ".webp": "image/webp"}
            mime = mime_map.get(ext, "image/png")

            response = self.openai_client.chat.completions.create(
                model=self.chat_model,
                messages=[
                    {"role": "system", "content": "Extract ALL text and content from this image. Return clean text preserving structure. Include tables in markdown format. Do NOT add commentary."},
                    {"role": "user", "content": [
                        {"type": "text", "text": "Extract all content from this image:"},
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"}}
                    ]}
                ],
                temperature=0.1, max_tokens=16000,
            )
            content = response.choices[0].message.content
            if content and len(content.strip()) >= 10:
                print(f"[DocumentParser] GPT-4o vision extracted {len(content)} chars from {file_name}")
                return content
        except Exception as e:
            print(f"[DocumentParser] GPT-4o vision error for {file_name}: {e}")

        if self.azure_endpoint and self.azure_api_key:
            return self._parse_with_azure_di_sync(file_bytes, file_name, ext)
        return ""

    def _parse_pdf_with_gpt4o_sync(self, file_bytes: bytes, file_name: str) -> str:
        """Synchronous GPT-4o PDF parsing (page-to-image)."""
        if not self.openai_client:
            return ""
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            all_text = []
            max_pages = min(len(doc), 20)
            for page_num in range(max_pages):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                b64 = base64.b64encode(img_bytes).decode('utf-8')
                response = self.openai_client.chat.completions.create(
                    model=self.chat_model,
                    messages=[
                        {"role": "system", "content": "Extract ALL text from this PDF page image. Preserve structure, tables, and formatting. Return only the extracted content."},
                        {"role": "user", "content": [
                            {"type": "text", "text": f"Extract all text from page {page_num + 1}:"},
                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"}}
                        ]}
                    ],
                    temperature=0.1, max_tokens=16000,
                )
                page_text = response.choices[0].message.content
                if page_text:
                    all_text.append(f"[Page {page_num + 1}]\n{page_text}")
            doc.close()
            return "\n\n".join(all_text) if all_text else ""
        except Exception as e:
            print(f"[DocumentParser] GPT-4o PDF sync error for {file_name}: {e}")
            return ""

    def _parse_with_llamaparse_sync(self, file_bytes: bytes, file_name: str, extension: str) -> str:
        """Synchronous LlamaParse using httpx sync client."""
        try:
            import httpx as httpx_mod
            ext = extension.lstrip('.')
            print(f"[DocumentParser] Parsing {file_name} with LlamaParse (sync)")

            with httpx_mod.Client(timeout=120.0) as client:
                files = {"file": (file_name, file_bytes, self._get_mime_type(ext))}
                data = {"gpt4o_mode": "true"}
                headers = {"Authorization": f"Bearer {self.llama_api_key}"}

                upload_response = client.post(f"{LLAMAPARSE_API_URL}/upload", files=files, data=data, headers=headers)
                if upload_response.status_code != 200:
                    print(f"[DocumentParser] LlamaParse upload failed: {upload_response.status_code}")
                    return ""

                job_id = upload_response.json().get("id")
                if not job_id:
                    return ""

                import time
                for attempt in range(60):
                    status_response = client.get(f"{LLAMAPARSE_API_URL}/job/{job_id}", headers=headers)
                    if status_response.status_code != 200:
                        time.sleep(5)
                        continue
                    status_data = status_response.json()
                    status = status_data.get("status")
                    if status == "SUCCESS":
                        break
                    elif status == "ERROR":
                        return ""
                    time.sleep(5)
                else:
                    return ""

                result_response = client.get(f"{LLAMAPARSE_API_URL}/job/{job_id}/result/text", headers=headers)
                if result_response.status_code != 200:
                    result_response = client.get(f"{LLAMAPARSE_API_URL}/job/{job_id}/result/markdown", headers=headers)
                if result_response.status_code == 200:
                    try:
                        rj = result_response.json()
                        text = rj.get("text") or rj.get("markdown") or result_response.text if isinstance(rj, dict) else result_response.text
                    except Exception:
                        text = result_response.text
                    print(f"[DocumentParser] LlamaParse sync: {len(text)} chars from {file_name}")
                    return text
                return ""
        except Exception as e:
            print(f"[DocumentParser] LlamaParse sync error: {e}")
            return ""

    def _parse_with_azure_di_sync(self, file_bytes: bytes, file_name: str, extension: str) -> str:
        """Synchronous Azure Document Intelligence."""
        try:
            import httpx as httpx_mod
            import time
            analyze_url = f"{self.azure_endpoint}/documentintelligence/documentModels/prebuilt-layout:analyze?api-version=2024-11-30"
            headers = {
                "Ocp-Apim-Subscription-Key": self.azure_api_key,
                "Content-Type": self._get_mime_type(extension.lstrip('.'))
            }
            with httpx_mod.Client(timeout=120.0) as client:
                response = client.post(analyze_url, content=file_bytes, headers=headers)
                if response.status_code != 202:
                    return ""
                operation_location = response.headers.get("Operation-Location")
                if not operation_location:
                    return ""
                poll_headers = {"Ocp-Apim-Subscription-Key": self.azure_api_key}
                for _ in range(60):
                    time.sleep(2)
                    status_response = client.get(operation_location, headers=poll_headers)
                    if status_response.status_code != 200:
                        continue
                    result = status_response.json()
                    status = result.get("status")
                    if status == "succeeded":
                        analyze_result = result.get("analyzeResult", {})
                        content = analyze_result.get("content", "")
                        if content:
                            return content
                        paragraphs = analyze_result.get("paragraphs", [])
                        if paragraphs:
                            return "\n\n".join([p.get("content", "") for p in paragraphs])
                        return ""
                    elif status == "failed":
                        return ""
                return ""
        except Exception as e:
            print(f"[DocumentParser] Azure DI sync error: {e}")
            return ""

    # ─── Utilities ──────────────────────────────────────────────────

    def _get_mime_type(self, extension: str) -> str:
        ext = extension.lower().lstrip('.')
        mime_types = {
            "pdf": "application/pdf",
            "doc": "application/msword",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "ppt": "application/vnd.ms-powerpoint",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "xls": "application/vnd.ms-excel",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "txt": "text/plain", "md": "text/markdown", "csv": "text/csv",
            "json": "application/json", "xml": "application/xml",
            "html": "text/html", "htm": "text/html", "rtf": "application/rtf",
            "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "gif": "image/gif", "bmp": "image/bmp", "tiff": "image/tiff", "tif": "image/tiff",
        }
        return mime_types.get(ext, "application/octet-stream")


# Singleton instance
_parser_instance: Optional[DocumentParser] = None


def get_document_parser(force_new: bool = False) -> DocumentParser:
    global _parser_instance
    if _parser_instance is None or force_new:
        _parser_instance = DocumentParser()
    return _parser_instance


def reset_parser():
    global _parser_instance
    _parser_instance = None


def parse_document_sync(file_bytes: bytes, file_name: str, extension: str) -> str:
    """Synchronous wrapper for parsing documents."""
    parser = get_document_parser()
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as executor:
                future = executor.submit(
                    asyncio.run,
                    parser.parse_bytes(file_bytes, file_name, extension)
                )
                return future.result()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    return loop.run_until_complete(
        parser.parse_bytes(file_bytes, file_name, extension)
    )
