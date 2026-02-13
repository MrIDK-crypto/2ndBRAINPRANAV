"""
Training PPT Generator
Generates training presentations using a template PPT as the base style
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import json
import copy

# Template path
TEMPLATE_PATH = "/Users/rishitjain/Downloads/BEAT x UCLA Final Presentation.pptx"
OUTPUT_DIR = "/Users/rishitjain/Downloads/knowledgevault_backend/training_generator/output"

# UCLA Health Training Content
UCLA_HEALTH_TRAINING = {
    "title": "UCLA Health NICU Project",
    "subtitle": "Training Module: FDU Relocation Business Plan",
    "sections": [
        {
            "type": "title",
            "title": "UCLA Health NICU Project",
            "subtitle": "Training Module Overview"
        },
        {
            "type": "section_header",
            "title": "PROJECT OVERVIEW"
        },
        {
            "type": "content",
            "title": "Executive Summary",
            "content": [
                "Project assesses two service expansion options to replace the existing FDU",
                "Option 1: NICU Step-Down Unit",
                "Option 2: OB-ED and Triage Unit",
                "Initial investment required: $1.42M",
                "Implementation timeline: 18-24 months"
            ]
        },
        {
            "type": "content",
            "title": "Problem Statement",
            "content": [
                "UCLA Health faces significant over-capacity in NICU and L&D beds",
                "Impact on patient safety, care quality, and satisfaction",
                "Healthcare worker morale affected",
                "44/420 transfer NICU patients turned away (2019-2024)",
                "288/420 transfer PICU patients turned away",
                "Lost opportunity cost: $632,553 annually"
            ]
        },
        {
            "type": "section_header",
            "title": "CASE BACKGROUND"
        },
        {
            "type": "content",
            "title": "Current State Analysis",
            "content": [
                "UCLA Health equipped with 20 NICU beds and 5 L&D beds",
                "Turning away critical NICU and PICU patients",
                "Prolonged stays worsen health outcomes",
                "L&D beds being used for patient triage",
                "Average NICU daily cost: $8,393",
                "Total addressable market (TAM): $812M for LA NICU services"
            ]
        },
        {
            "type": "content",
            "title": "Competitor Analysis",
            "content": [
                "Cedars-Sinai: Major competitor in LA market",
                "Children's Hospital Los Angeles (CHLA): Key pediatric competitor",
                "Boston Children's Hospital: #1 Regional Ranking (New England)",
                "UCLA Mattel Children's Hospital: Currently unranked in Pacific region",
                "Most Honor Roll hospitals have NICU Step-Down units"
            ]
        },
        {
            "type": "section_header",
            "title": "FINANCIAL ANALYSIS"
        },
        {
            "type": "content",
            "title": "NICU Step-Down Financial Model",
            "content": [
                "Initial Investment: $2,441,300",
                "Year 1 Revenue: $8,987,557 | ROI: 14%",
                "Year 2 Revenue: $9,442,328 | ROI: 120%",
                "Year 3 Revenue: $9,920,110 | ROI: 125%",
                "NPV (7% discount rate): $6,855,853",
                "Annual patients: 338 (Year 1) to 358 (Year 3)"
            ]
        },
        {
            "type": "content",
            "title": "OB-ED Financial Model",
            "content": [
                "Initial Investment: $2,290,000",
                "Year 1 Revenue: $5,751,801 | ROI: -74%",
                "Year 2 Revenue: $6,007,056 | ROI: 25%",
                "Year 3 Revenue: $6,424,776 | ROI: 28%",
                "Annual patients: 3,500 to 3,825",
                "Revenue per patient: ~$1,643"
            ]
        },
        {
            "type": "content",
            "title": "Market Size Analysis",
            "content": [
                "LA Population: 3,820,000",
                "Total Annual Births: 96,230",
                "NICU utilization rate: 10%",
                "NICU Step-Down Utilization: 60%",
                "TAM: $812M | SAM: $188M | SOM: $47M",
                "New SOM with Step-Down: $56.5M"
            ]
        },
        {
            "type": "section_header",
            "title": "RECOMMENDATION"
        },
        {
            "type": "content",
            "title": "Final Recommendation",
            "content": [
                "RECOMMENDED: Creation of the OB-ED Unit",
                "Addresses overcapacity issues in L&D",
                "Improves patient care transition",
                "Better long-term benefits vs NICU Step-Down",
                "Key components: Dedicated OBED unit with 5 beds",
                "Specialized triage, obstetric emergencies, post-treatment areas"
            ]
        },
        {
            "type": "content",
            "title": "Expected Outcomes",
            "content": [
                "Enhanced Patient Care: Faster obstetric emergency response",
                "Improved maternal and neonatal outcomes",
                "Operational Efficiency: Better L&D resource utilization",
                "Reduced turnover time in L&D units",
                "Reduced L&D Over-Capacity: Shifts emergency cases to OBED",
                "Long-term cost savings and increased patient throughput"
            ]
        },
        {
            "type": "section_header",
            "title": "RISKS & MITIGATION"
        },
        {
            "type": "content",
            "title": "NICU Step-Down Risks",
            "content": [
                "Risk 1: Patient flow management challenges",
                "Solution: Clear protocols and effective communication",
                "Risk 2: Parent anxiety during transition",
                "Solution: Comprehensive family education and counseling",
                "Risk 3: Staffing and training requirements",
                "Solution: Phased recruitment and ongoing training programs"
            ]
        },
        {
            "type": "content",
            "title": "OB-ED Risks",
            "content": [
                "Risk 1: Regulatory and compliance issues",
                "Solution: Work with departments, regular policy reviews",
                "Risk 2: Resistance from community physicians",
                "Solution: Engage physicians in planning, highlight benefits",
                "Risk 3: Nurse turnover",
                "Solution: Competitive salaries, professional development"
            ]
        },
        {
            "type": "section_header",
            "title": "SUMMARY"
        },
        {
            "type": "content",
            "title": "Key Takeaways",
            "content": [
                "UCLA Health faces critical capacity challenges in NICU/L&D",
                "Two viable options evaluated: NICU Step-Down and OB-ED",
                "OB-ED recommended for better long-term outcomes",
                "Initial investment: $2.29M with ROI turning positive Year 2",
                "Implementation: 18-24 months",
                "Success metrics: Patient throughput, satisfaction, ROI"
            ]
        },
        {
            "type": "title",
            "title": "Questions?",
            "subtitle": "UCLA Health NICU Project Training"
        }
    ]
}


def create_training_ppt(content: dict, template_path: str, output_path: str):
    """Create a training PPT using the template as base"""

    # Load the template
    prs = Presentation(template_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the key layouts we'll use
    # Layout 1: Section Header (for section dividers)
    # Layout 0: Title and Content (for content slides)

    # We'll clone slides from the template to maintain the exact styling
    # Key slides to use as templates:
    # Slide 1 (index 0): Title slide
    # Slide 4 (index 3): Section header (CASE BACKGROUND style)
    # Slide 5 (index 4): Content with title and bullets

    # Create a new presentation using the template layouts
    new_prs = Presentation(template_path)

    # Remove all existing slides
    while len(new_prs.slides) > 0:
        rId = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(rId)
        del new_prs.slides._sldIdLst[0]

    # Get layouts
    section_header_layout = None
    content_layout = None
    title_layout = None

    for layout in new_prs.slide_layouts:
        if "Section Header" in layout.name:
            section_header_layout = layout
        elif "UConsulting - Meet the Team" in layout.name:
            title_layout = layout
        elif "Title and Content" in layout.name and content_layout is None:
            content_layout = layout

    # Fallbacks
    if section_header_layout is None:
        section_header_layout = new_prs.slide_layouts[2]  # Section Header
    if content_layout is None:
        content_layout = new_prs.slide_layouts[0]  # Title and Content
    if title_layout is None:
        title_layout = new_prs.slide_layouts[1]  # UConsulting Meet the Team

    # Colors from template (UCLA Blue theme)
    UCLA_BLUE = RGBColor(0x2D, 0x68, 0xC4)  # Primary blue
    DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)  # Dark navy
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def add_title_slide(title_text: str, subtitle_text: str = ""):
        """Add a title/section slide"""
        slide = new_prs.slides.add_slide(title_layout)

        # Add decorative shapes (matching template style)
        # Left blue bar
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.5), slide_height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = UCLA_BLUE
        left_bar.line.fill.background()

        # Add title text box
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(24)
            p2.font.color.rgb = UCLA_BLUE
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)

        return slide

    def add_section_header(title_text: str):
        """Add a section header slide"""
        slide = new_prs.slides.add_slide(section_header_layout)

        # Add left blue bar
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.5), slide_height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = UCLA_BLUE
        left_bar.line.fill.background()

        # Add diagonal accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM,
            Inches(0.5), Inches(2),
            Inches(4), Inches(3)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = UCLA_BLUE
        accent.fill.fore_color.brightness = 0.3
        accent.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3),
            Inches(10), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        return slide

    def add_content_slide(title_text: str, bullets: list):
        """Add a content slide with title and bullets"""
        slide = new_prs.slides.add_slide(content_layout)

        # Add left blue bar
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), slide_height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = UCLA_BLUE
        left_bar.line.fill.background()

        # Add top accent line
        top_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(0),
            Inches(13), Inches(0.05)
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = UCLA_BLUE
        top_line.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(11), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Add horizontal line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.2),
            Inches(11.5), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = UCLA_BLUE
        line.line.fill.background()

        # Add bullet points
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(11.5), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_before = Pt(12)
            p.space_after = Pt(6)
            p.level = 0

        # Add slide number
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7),
            Inches(0.5), Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(len(new_prs.slides))
        p.font.size = Pt(12)
        p.font.color.rgb = UCLA_BLUE
        p.alignment = PP_ALIGN.RIGHT

        return slide

    # Generate slides from content
    for section in content["sections"]:
        if section["type"] == "title":
            add_title_slide(section["title"], section.get("subtitle", ""))
        elif section["type"] == "section_header":
            add_section_header(section["title"])
        elif section["type"] == "content":
            add_content_slide(section["title"], section["content"])

    # Save the presentation
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    new_prs.save(output_path)
    print(f"Training PPT saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    output_path = os.path.join(OUTPUT_DIR, "UCLA_Health_Training.pptx")
    create_training_ppt(UCLA_HEALTH_TRAINING, TEMPLATE_PATH, output_path)
    print(f"\nGenerated {len(UCLA_HEALTH_TRAINING['sections'])} slides")
